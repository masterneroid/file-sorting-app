import sys
import os
import shutil
import json
import time
import logging
import multiprocessing
import concurrent.futures
import threading
from datetime import datetime
from typing import Dict, List, Tuple, Optional, Any
from collections import Counter
import re

# psutil'u deneyerek import et
try:
    import psutil
    PSUTIL_AVAILABLE = True
except ImportError:
    PSUTIL_AVAILABLE = False
    print("Uyarı: psutil kütüphanesi kurulu değil. Sistem optimizasyonları devre dışı.")
    print("Kurmak için: pip install psutil")

# Doküman işleme kütüphaneleri
try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Uyarı: pypdf kütüphanesi kurulu değil. PDF desteği devre dışı.")
    print("Kurmak için: pip install pypdf")

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Uyarı: python-docx kütüphanesi kurulu değil. DOCX desteği devre dışı.")
    print("Kurmak için: pip install python-docx")

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Uyarı: python-pptx kütüphanesi kurulu değil. PPTX desteği devre dışı.")
    print("Kurmak için: pip install python-pptx")

try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("Uyarı: openpyxl kütüphanesi kurulu değil. Excel (.xlsx) desteği devre dışı.")
    print("Kurmak için: pip install openpyxl")

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False
    print("Uyarı: xlrd kütüphanesi kurulu değil. Eski Excel (.xls) desteği devre dışı.")
    print("Kurmak için: pip install xlrd")

# AI/ML kütüphaneleri
try:
    import torch
    from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM
    TRANSFORMERS_AVAILABLE = True
except ImportError:
    TRANSFORMERS_AVAILABLE = False
    print("Uyarı: transformers kütüphanesi kurulu değil. AI analizi devre dışı.")
    print("Kurmak için: pip install transformers torch")

from PySide6.QtWidgets import (
    QApplication, QMainWindow, QPushButton, QVBoxLayout, QWidget, 
    QFileDialog, QLabel, QMessageBox, QHBoxLayout, QCheckBox, 
    QProgressBar, QGroupBox, QTextEdit, QTabWidget, QComboBox,
    QSpinBox, QTableWidget, QTableWidgetItem, QHeaderView, QSlider,
    QSplitter, QTreeWidget, QTreeWidgetItem, QLineEdit
)
from PySide6.QtCore import Qt, QUrl, Signal, QObject, QThread, QTimer
from PySide6.QtGui import QDesktopServices, QFont, QIcon, QTextCursor

# ==================== DOCUMENT ANALYZER ====================

class DocumentAnalyzer:
    """Doküman analizi ve içerik bazlı sınıflandırma"""
    
    def __init__(self, config):
        self.config = config
        self.ai_pipeline = None
        self.initialized = False
        
        # Türkçe anahtar kelime kategorileri
        self.categories = {
            "Finans": ["bütçe", "fatura", "ödeme", "bank", "para", "hesap", "mali", "finans", 
                      "kredi", "borç", "yatırım", "borsa", "döviz", "vergi", "maaş"],
            "Eğitim": ["ders", "ödev", "proje", "sınav", "okul", "üniversite", "eğitim", "öğrenci",
                      "not", "sınıf", "kurs", "seminer", "akademik", "tez", "ders notu"],
            "İş": ["rapor", "toplantı", "proje", "sunum", "iş", "şirket", "yönetim", "strateji",
                  "plan", "çalışan", "müdür", "müşteri", "satış", "pazarlama", "insan kaynakları"],
            "Teknik": ["kod", "yazılım", "donanım", "teknik", "sistem", "network", "server", "database",
                      "program", "algoritma", "yapay zeka", "machine learning", "veri", "analiz"],
            "Sağlık": ["sağlık", "hasta", "tedavi", "rapor", "ilaç", "doktor", "hastane", "muayene",
                      "tahlil", "reçete", "ameliyat", "tedavi", "psikoloji", "terapi"],
            "Hukuk": ["sözleşme", "kanun", "yasa", "hukuk", "dava", "avukat", "mahkeme", "anlaşma",
                     "taraflar", "madde", "yargı", "ceza", "hüküm", "temyiz"],
            "Kişisel": ["cv", "özgeçmiş", "mektup", "kişisel", "iletişim", "aile", "arkadaş", "ev",
                       "tatil", "gezi", "günlük", "anı", "fotoğraf", "video"],
            "Araştırma": ["araştırma", "makale", "tez", "bilim", "akademik", "yayın", "doktora",
                         "literatür", "deney", "sonuç", "hipotez", "bulgu", "analiz"],
            "Tasarım": ["tasarım", "çizim", "grafik", "resim", "şekil", "layout", "ui", "ux",
                       "renk", "font", "illustrasyon", "mockup", "prototip"],
            "Yönetim": ["plan", "strateji", "hedef", "performans", "kalite", "süreç", "proje yönetimi",
                       "risk", "bütçe", "kpi", "rapor", "analiz", "karar"]
        }
        
        # Türkçe stop words
        self.turkish_stopwords = {
            'acaba', 'ama', 'aslında', 'az', 'bazı', 'belki', 'biri', 'birkaç', 'birşey', 'biz',
            'bu', 'çok', 'çünkü', 'da', 'daha', 'de', 'defa', 'diye', 'eğer', 'en', 'gibi', 'hem',
            'hep', 'hepsi', 'her', 'hiç', 'için', 'ile', 'ise', 'kez', 'ki', 'kim', 'mı', 'mu',
            'mü', 'nasıl', 'ne', 'neden', 'nerde', 'nerede', 'nereye', 'niçin', 'niye', 'o', 'sanki',
            'şey', 'siz', 'şu', 'tüm', 've', 'veya', 'ya', 'yani'
        }
        
        # AI modelini başlat
        self._initialize_ai()
    
    def _initialize_ai(self):
        """AI modelini başlat"""
        try:
            if TRANSFORMERS_AVAILABLE and self.config.get("ai_enabled", True):
                model_name = "facebook/bart-large-mnli"  # Zero-shot classification için iyi
                
                self.ai_pipeline = pipeline(
                    "zero-shot-classification",
                    model=model_name,
                    device=-1  # CPU kullan
                )
                self.initialized = True
                logging.info("AI model başarıyla yüklendi")
            else:
                logging.info("AI model kullanılamıyor, anahtar kelime analizi kullanılacak")
        except Exception as e:
            logging.error(f"AI model başlatma hatası: {e}")
            self.initialized = False
    
    def extract_text_from_pdf(self, filepath):
        """PDF'den metin çıkar - Geliştirilmiş versiyon"""
        if not PDF_AVAILABLE:
            print(f"PDF desteği kapalı: {filepath}")
            return ""
        
        try:
            logging.info(f"PDF okunuyor: {filepath}")
            
            # PDF'i aç
            reader = PdfReader(filepath)
            text = ""
            total_pages = len(reader.pages)
            
            logging.info(f"  Toplam sayfa: {total_pages}")
            
            # Her sayfayı oku
            for i, page in enumerate(reader.pages):
                try:
                    # Metni çıkar
                    page_text = page.extract_text()
                    
                    if page_text and page_text.strip():
                        # Temizle ve ekle
                        cleaned_text = page_text.strip()
                        text += cleaned_text + "\n\n"
                        
                        # Büyük PDF'ler için ilerleme bilgisi
                        if total_pages > 20 and (i + 1) % 10 == 0:
                            logging.info(f"  Sayfa {i + 1}/{total_pages} okundu")
                    else:
                        # Boş veya metin içermeyen sayfa
                        logging.debug(f"  Sayfa {i + 1} boş veya metin içermiyor")
                        
                except Exception as page_error:
                    logging.warning(f"  Sayfa {i + 1} okuma hatası: {page_error}")
                    continue
            
            logging.info(f"PDF okuma tamamlandı: {len(text)} karakter çıkarıldı")
            return text.strip()
            
        except Exception as e:
            logging.error(f"PDF okuma hatası {filepath}: {e}")
            return ""
    
    def extract_text_from_docx(self, filepath):
        """DOCX'ten metin çıkar"""
        if not DOCX_AVAILABLE:
            return ""
        
        try:
            doc = docx.Document(filepath)
            text = ""
            
            # Paragraflar
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Tablolar
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            return text.strip()
        except Exception as e:
            logging.error(f"DOCX okuma hatası {filepath}: {e}")
            return ""
    
    def extract_text_from_pptx(self, filepath):
        """PPTX'ten metin çıkar"""
        if not PPTX_AVAILABLE:
            return ""
        
        try:
            prs = pptx.Presentation(filepath)
            text = ""
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text += "\n".join(slide_text) + "\n---\n"
            
            return text.strip()
        except Exception as e:
            logging.error(f"PPTX okuma hatası {filepath}: {e}")
            return ""
    
    def extract_text_from_excel(self, filepath):
        """Excel'den metin çıkar"""
        ext = os.path.splitext(filepath)[1].lower()
        
        if ext == '.xlsx':
            if not EXCEL_AVAILABLE:
                return ""
            try:
                wb = load_workbook(filepath, data_only=True, read_only=True)
                text = ""
                
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_text = []
                    
                    for row in ws.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                cell_text = str(cell).strip()
                                if cell_text:
                                    row_text.append(cell_text)
                        
                        if row_text:
                            sheet_text.append(" | ".join(row_text))
                    
                    if sheet_text:
                        text += f"### {sheet_name} ###\n" + "\n".join(sheet_text) + "\n\n"
                
                return text.strip()
            except Exception as e:
                logging.error(f"Excel (.xlsx) okuma hatası {filepath}: {e}")
                return ""
        
        elif ext == '.xls':
            if not XLRD_AVAILABLE:
                return ""
            try:
                import xlrd
                workbook = xlrd.open_workbook(filepath)
                text = ""
                
                for sheet in workbook.sheets():
                    sheet_text = []
                    
                    for row_idx in range(sheet.nrows):
                        row_text = []
                        for col_idx in range(sheet.ncols):
                            cell = sheet.cell(row_idx, col_idx)
                            if cell.value:
                                cell_text = str(cell.value).strip()
                                if cell_text:
                                    row_text.append(cell_text)
                        
                        if row_text:
                            sheet_text.append(" | ".join(row_text))
                    
                    if sheet_text:
                        text += f"### {sheet.name} ###\n" + "\n".join(sheet_text) + "\n\n"
                
                return text.strip()
            except Exception as e:
                logging.error(f"Excel (.xls) okuma hatası {filepath}: {e}")
                return ""
        
        return ""
    
    def extract_text_from_txt(self, filepath):
        """TXT dosyasından metin çıkar"""
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read().strip()
        except UnicodeDecodeError:
            try:
                with open(filepath, 'r', encoding='latin-1') as f:
                    return f.read().strip()
            except Exception as e:
                logging.error(f"TXT okuma hatası {filepath}: {e}")
                return ""
        except Exception as e:
            logging.error(f"TXT okuma hatası {filepath}: {e}")
            return ""
    
    def extract_text(self, filepath):
        """Dosya tipine göre metin çıkar"""
        if not os.path.exists(filepath):
            logging.error(f"Dosya bulunamadı: {filepath}")
            return ""
        
        ext = os.path.splitext(filepath)[1].lower()
        
        logging.info(f"Metin çıkarılıyor: {filepath} ({ext})")
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(filepath)
        elif ext == '.docx':
            return self.extract_text_from_docx(filepath)
        elif ext == '.pptx':
            return self.extract_text_from_pptx(filepath)
        elif ext in ['.xlsx', '.xls']:
            return self.extract_text_from_excel(filepath)
        elif ext == '.txt':
            return self.extract_text_from_txt(filepath)
        elif ext == '.md':
            return self.extract_text_from_txt(filepath)
        else:
            logging.warning(f"Desteklenmeyen dosya türü: {ext}")
            return ""
    
    def analyze_content(self, text, filename=""):
        """Metin içeriğini analiz et ve kategori öner"""
        if not text or len(text.strip()) < 20:
            return "Diğer_Belgeler", [], "İçerik yetersiz"
        
        try:
            # 1. Anahtar kelime analizi
            keyword_categories = self._keyword_analysis(text)
            
            # 2. AI analizi (eğer aktifse)
            ai_categories = []
            if self.initialized and len(text) > 100:
                ai_categories = self._ai_analysis(text)
            
            # 3. Tüm kategorileri birleştir
            all_categories = keyword_categories + ai_categories
            
            if not all_categories:
                return "Diğer_Belgeler", [], self._create_summary(text)
            
            # 4. En olası kategoriyi bul
            category_counter = Counter(all_categories)
            main_category = category_counter.most_common(1)[0][0]
            
            # 5. İlgili kategoriler
            related = [cat for cat, count in category_counter.most_common(3) if cat != main_category]
            
            # 6. Özet oluştur
            summary = self._create_summary(text)
            
            return main_category, related, summary
            
        except Exception as e:
            logging.error(f"İçerik analiz hatası: {e}")
            return "Diğer_Belgeler", [], "Analiz hatası"
    
    def _keyword_analysis(self, text):
        """Anahtar kelime analizi ile kategori tespiti"""
        text_lower = text.lower()
        
        # Noktalama işaretlerini kaldır
        text_clean = re.sub(r'[^\w\s]', ' ', text_lower)
        words = text_clean.split()
        
        # Stop words'leri filtrele
        words_filtered = [w for w in words if w not in self.turkish_stopwords and len(w) > 2]
        
        detected = []
        word_counter = Counter(words_filtered)
        
        for category, keywords in self.categories.items():
            score = 0
            for keyword in keywords:
                # Anahtar kelimenin metinde geçme sayısı
                score += word_counter.get(keyword, 0)
            
            # Eşik değeri
            if score >= 2:
                detected.append(category)
        
        return detected
    
    def _ai_analysis(self, text):
        """AI ile kategori analizi"""
        try:
            # Kategori etiketlerini hazırla
            candidate_labels = list(self.categories.keys())
            
            # Metni kısalt (model sınırları için)
            text_sample = text[:500]
            
            result = self.ai_pipeline(
                text_sample,
                candidate_labels=candidate_labels,
                multi_label=True
            )
            
            # Güven skoru yüksek olan kategorileri al
            categories = []
            for label, score in zip(result['labels'], result['scores']):
                if score > 0.3:  # %30 eşik değeri
                    categories.append(label)
            
            return categories
            
        except Exception as e:
            logging.error(f"AI analiz hatası: {e}")
            return []
    
    def _create_summary(self, text, max_length=150):
        """Metin özeti oluştur"""
        if not text:
            return ""
        
        try:
            # İlk birkaç cümleyi al
            sentences = re.split(r'[.!?]+', text)
            valid_sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
            
            if not valid_sentences:
                # Uzun bir kelime dizisi varsa ilk 100 karakter
                return text[:max_length] + '...' if len(text) > max_length else text
            
            summary = ""
            for sentence in valid_sentences[:3]:
                if len(summary) + len(sentence) < max_length:
                    summary += sentence + '. '
                else:
                    break
            
            summary = summary.strip()
            if len(summary) > max_length:
                summary = summary[:max_length] + '...'
            
            return summary
            
        except Exception as e:
            logging.error(f"Özet oluşturma hatası: {e}")
            return text[:100] + '...'
    
    def extract_metadata(self, filepath, text):
        """Dosya metadata'sını çıkar"""
        try:
            metadata = {
                "dosya_adi": os.path.basename(filepath),
                "dosya_yolu": filepath,
                "dosya_boyutu": os.path.getsize(filepath),
                "son_duzenleme": datetime.fromtimestamp(os.path.getmtime(filepath)).isoformat(),
                "karakter_sayisi": len(text),
                "kelime_sayisi": len(text.split()),
                "satir_sayisi": text.count('\n') + 1
            }
            
            # Anahtar kelimeler
            words = re.findall(r'\b\w{4,}\b', text.lower())
            filtered_words = [w for w in words if w not in self.turkish_stopwords]
            top_keywords = Counter(filtered_words).most_common(5)
            
            metadata["anahtar_kelimeler"] = [word for word, count in top_keywords]
            
            return metadata
            
        except Exception as e:
            logging.error(f"Metadata çıkarma hatası: {e}")
            return {}
        # ... extract_metadata fonksiyonunun bittiği yerden itibaren ...
    
    def create_category_structure(self, base_path, category):
        """Kategori klasörünü fiziksel olarak oluşturur ve yolu döndürür."""
        import os # Hata almamak için fonksiyon içinde de tanımlayabiliriz
        target_dir = os.path.join(base_path, category)
        if not os.path.exists(target_dir):
            os.makedirs(target_dir, exist_ok=True)
        return target_dir

    def organize_file(self, source_file, category, output_root):
        """Dosyayı fiziksel olarak kategorisine taşır."""
        import os
        import shutil
        import time
        try:
            target_dir = self.create_category_structure(output_root, category)
            filename = os.path.basename(source_file)
            dest_path = os.path.join(target_dir, filename)
            
            # Dosya zaten varsa ismini değiştir (Üzerine yazmamak için)
            if os.path.exists(dest_path):
                base, ext = os.path.splitext(filename)
                dest_path = os.path.join(target_dir, f"{base}_{int(time.time())}{ext}")
            
            shutil.copy2(source_file, dest_path)
            return True, dest_path
        except Exception as e:
            return False, str(e)
      
        
        

# ==================== THEME MANAGER ====================

class ThemeManager:
    """Tema yöneticisi sınıfı"""
    
    THEMES = {
        "light": {
           "window_bg": "#FFFFFF",
        "widget_bg": "#F8F9FA",
        "widget_bg_alt": "#FFFFFF",
        "text_primary": "#212529",
        "text_secondary": "#6C757D",
        "text_inverse": "#FFFFFF",
        "border": "#DEE2E6",
        "border_light": "#E9ECEF",
        "accent_primary": "#007BFF",
        "accent_primary_hover": "#0056B3",
        "accent_success": "#28A745",
        "accent_success_hover": "#1E7E34",
        "accent_warning": "#FFC107",
        "accent_warning_hover": "#E0A800",
        "accent_danger": "#DC3545",
        "accent_danger_hover": "#BD2130",
        "accent_info": "#17A2B8",
        "accent_info_hover": "#138496",
        "tab_bg": "#E9ECEF",
        "tab_bg_selected": "#FFFFFF",
        "tab_border": "#DEE2E6",
        "header_bg": "#F8F9FA",
        "header_text": "#495057",
        "progress_bg": "#E9ECEF",
        "progress_chunk": "#28A745",
        "groupbox_bg": "#FFFFFF",
        "groupbox_border": "#DEE2E6",
        "button_disabled": "#E9ECEF",
        "button_disabled_text": "#6C757D"
        },
        "dark": {
            "window_bg": "#1a1a1a",
            "widget_bg": "#2d2d2d",
            "widget_bg_alt": "#3d3d3d",
            "text_primary": "#f8f9fa",
            "text_secondary": "#adb5bd",
            "text_inverse": "#212529",
            "border": "#495057",
            "border_light": "#6c757d",
            "accent_primary": "#3498db",
            "accent_primary_hover": "#2980b9",
            "accent_success": "#2ecc71",
            "accent_success_hover": "#27ae60",
            "accent_warning": "#f39c12",
            "accent_warning_hover": "#e67e22",
            "accent_danger": "#e74c3c",
            "accent_danger_hover": "#c0392b",
            "accent_info": "#17a2b8",
            "accent_info_hover": "#138496",
            "tab_bg": "#2d2d2d",
            "tab_bg_selected": "#3d3d3d",
            "tab_border": "#495057",
            "header_bg": "#3d3d3d",
            "header_text": "#f8f9fa",
            "progress_bg": "#3d3d3d",
            "progress_chunk": "#2ecc71",
            "groupbox_bg": "#2d2d2d",
            "groupbox_border": "#495057",
            "button_disabled": "#6c757d",
            "button_disabled_text": "#adb5bd"
        }
    }
    
    @staticmethod
    def get_theme(theme_name: str) -> dict:
        """Tema adına göre tema sözlüğünü getir"""
        return ThemeManager.THEMES.get(theme_name, ThemeManager.THEMES["light"])
    
    @staticmethod
    def generate_stylesheet(theme: dict) -> str:
        """Tema sözlüğünden stylesheet oluştur"""
        return f"""
            /* Ana Pencere */
            QMainWindow {{
                background-color: {theme['window_bg']};
            }}
            
            /* Tüm Widget'lar */
            QWidget {{
                background-color: transparent;
                color: {theme['text_primary']};
                font-family: 'Segoe UI', Arial, sans-serif;
                font-size: 13px;
            }}
            
            /* Grup Kutuları */
            QGroupBox {{
                font-weight: bold;
                border: 2px solid {theme['groupbox_border']};
                border-radius: 8px;
                margin-top: 10px;
                padding-top: 12px;
                background-color: {theme['groupbox_bg']};
            }}
            
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 10px;
                padding: 0 8px 0 8px;
                color: {theme['text_primary']};
            }}
            
            /* Etiketler */
            QLabel {{
                color: {theme['text_primary']};
            }}
            
            QLabel[cssClass="title"] {{
                font-size: 20px;
                font-weight: bold;
                color: {theme['accent_primary']};
            }}
            
            QLabel[cssClass="subtitle"] {{
                color: {theme['text_secondary']};
                font-size: 12px;
            }}
            
            /* Metin Editleri */
            QTextEdit {{
                border: 1px solid {theme['border']};
                border-radius: 6px;
                padding: 8px;
                background-color: {theme['widget_bg']};
                color: {theme['text_primary']};
                selection-background-color: {theme['accent_primary']};
                selection-color: {theme['text_inverse']};
            }}
            
            QTextEdit:focus {{
                border: 2px solid {theme['accent_primary']};
            }}
            
            /* Butonlar - Genel */
            QPushButton {{
                background-color: {theme['accent_primary']};
                color: {theme['text_inverse']};
                font-weight: bold;
                border: none;
                border-radius: 6px;
                padding: 10px 16px;
                font-size: 13px;
                min-height: 40px;
            }}
            
            QPushButton:hover {{
                background-color: {theme['accent_primary_hover']};
            }}
            
            QPushButton:pressed {{
                background-color: {theme['accent_primary_hover']};
                padding-top: 11px;
                padding-bottom: 9px;
            }}
            
            QPushButton:disabled {{
                background-color: {theme['button_disabled']};
                color: {theme['button_disabled_text']};
            }}
            
            /* Özel Butonlar */
            QPushButton[cssClass="success"] {{
                background-color: {theme['accent_success']};
            }}
            
            QPushButton[cssClass="success"]:hover {{
                background-color: {theme['accent_success_hover']};
            }}
            
            QPushButton[cssClass="danger"] {{
                background-color: {theme['accent_danger']};
            }}
            
            QPushButton[cssClass="danger"]:hover {{
                background-color: {theme['accent_danger_hover']};
            }}
            
            QPushButton[cssClass="warning"] {{
                background-color: {theme['accent_warning']};
            }}
            
            QPushButton[cssClass="warning"]:hover {{
                background-color: {theme['accent_warning_hover']};
            }}
            
            QPushButton[cssClass="info"] {{
                background-color: {theme['accent_info']};
            }}
            
            QPushButton[cssClass="info"]:hover {{
                background-color: {theme['accent_info_hover']};
            }}
            
            /* CheckBox'lar */
            QCheckBox {{
                spacing: 8px;
                color: {theme['text_primary']};
            }}
            
            QCheckBox::indicator {{
                width: 18px;
                height: 18px;
                border: 2px solid {theme['border']};
                border-radius: 3px;
            }}
            
            QCheckBox::indicator:checked {{
                background-color: {theme['accent_primary']};
                border-color: {theme['accent_primary']};
                image: url(:/images/check.png);
            }}
            
            QCheckBox::indicator:hover {{
                border-color: {theme['accent_primary']};
            }}
            
            /* ComboBox'lar */
            QComboBox {{
                border: 1px solid {theme['border']};
                border-radius: 6px;
                padding: 6px;
                background-color: {theme['widget_bg']};
                color: {theme['text_primary']};
                min-width: 120px;
            }}
            
            QComboBox::drop-down {{
                border: none;
                width: 20px;
            }}
            
            QComboBox::down-arrow {{
                image: none;
                border-left: 1px solid {theme['border']};
                width: 10px;
                height: 10px;
            }}
            
            QComboBox QAbstractItemView {{
                background-color: {theme['widget_bg']};
                color: {theme['text_primary']};
                border: 1px solid {theme['border']};
                selection-background-color: {theme['accent_primary']};
                selection-color: {theme['text_inverse']};
            }}
            
            /* SpinBox'lar */
            QSpinBox {{
                border: 1px solid {theme['border']};
                border-radius: 6px;
                padding: 6px;
                background-color: {theme['widget_bg']};
                color: {theme['text_primary']};
            }}
            
            QSpinBox::up-button, QSpinBox::down-button {{
                width: 20px;
                border: 1px solid {theme['border']};
            }}
            
            /* Slider */
            QSlider::groove:horizontal {{
                border: 1px solid {theme['border']};
                height: 8px;
                background: {theme['widget_bg']};
                margin: 2px 0;
                border-radius: 4px;
            }}
            
            QSlider::handle:horizontal {{
                background: {theme['accent_primary']};
                border: 2px solid {theme['accent_primary_hover']};
                width: 18px;
                height: 18px;
                margin: -5px 0;
                border-radius: 9px;
            }}
            
            QSlider::sub-page:horizontal {{
                background: {theme['accent_primary']};
                border: 1px solid {theme['border']};
                height: 8px;
                border-radius: 4px;
            }}
            
            /* TabWidget */
            QTabWidget::pane {{
                border: 1px solid {theme['tab_border']};
                border-radius: 6px;
                background-color: {theme['widget_bg']};
                margin-top: -1px;
            }}
            
            QTabBar::tab {{
                background-color: {theme['tab_bg']};
                color: {theme['text_primary']};
                padding: 8px 16px;
                margin-right: 2px;
                border: 1px solid {theme['tab_border']};
                border-bottom: none;
                border-top-left-radius: 6px;
                border-top-right-radius: 6px;
            }}
            
            QTabBar::tab:selected {{
                background-color: {theme['tab_bg_selected']};
                border-bottom: 2px solid {theme['accent_primary']};
                font-weight: bold;
            }}
            
            QTabBar::tab:hover:!selected {{
                background-color: {theme['widget_bg_alt']};
            }}
            
            /* Tablolar */
            QTableWidget {{
                border: 1px solid {theme['border']};
                border-radius: 6px;
                background-color: {theme['widget_bg']};
                alternate-background-color: {theme['widget_bg_alt']};
                gridline-color: {theme['border_light']};
                color: {theme['text_primary']};
            }}
            
            QHeaderView::section {{
                background-color: {theme['header_bg']};
                color: {theme['header_text']};
                padding: 10px;
                border: none;
                border-bottom: 1px solid {theme['border']};
                font-weight: bold;
            }}
            
            QTableWidget::item {{
                padding: 8px;
            }}
            
            QTableWidget::item:selected {{
                background-color: {theme['accent_primary']};
                color: {theme['text_inverse']};
            }}
            
            /* TreeWidget */
            QTreeWidget {{
                border: 1px solid {theme['border']};
                border-radius: 6px;
                background-color: {theme['widget_bg']};
                color: {theme['text_primary']};
            }}
            
            QTreeWidget::item {{
                padding: 4px;
            }}
            
            QTreeWidget::item:selected {{
                background-color: {theme['accent_primary']};
                color: {theme['text_inverse']};
            }}
            
            /* Progress Bar */
            QProgressBar {{
                border: 1px solid {theme['border']};
                border-radius: 6px;
                text-align: center;
                background-color: {theme['progress_bg']};
                color: {theme['text_primary']};
                height: 24px;
            }}
            
            QProgressBar::chunk {{
                background-color: {theme['progress_chunk']};
                border-radius: 6px;
            }}
            
            /* Status Bar */
            QStatusBar {{
                background-color: {theme['widget_bg']};
                color: {theme['text_primary']};
                border-top: 1px solid {theme['border']};
                padding: 4px;
            }}
            
            /* Splitter */
            QSplitter::handle {{
                background-color: {theme['border']};
            }}
            
            /* ScrollBar */
            QScrollBar:vertical {{
                background-color: {theme['widget_bg_alt']};
                width: 12px;
                border-radius: 6px;
            }}
            
            QScrollBar::handle:vertical {{
                background-color: {theme['border']};
                border-radius: 6px;
                min-height: 20px;
            }}
            
            QScrollBar::handle:vertical:hover {{
                background-color: {theme['accent_primary']};
            }}
            
            QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {{
                height: 0px;
            }}
        """

# ==================== CONFIGURATION CLASSES ====================

class ConfigManager:
    """Yapılandırma yönetimi için sınıf"""
    
    def __init__(self, config_file="ayarlar.json"):
        self.config_file = config_file
        self.default_config = {
            "kategoriler": {
                "Görseller": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg", ".webp", ".tiff", ".ico"],
                "Belgeler": [".pdf", ".docx", ".doc", ".txt", ".rtf", ".odt", ".md", ".tex"],
                "Ofis_Dosyaları": [".xlsx", ".xls", ".pptx", ".ppt", ".csv", ".ods", ".odp"],
                "Videolar": [".mp4", ".mkv", ".mov", ".avi", ".wmv", ".flv", ".webm", ".m4v", ".3gp"],
                "Müzik": [".mp3", ".wav", ".flac", ".m4a", ".aac", ".ogg", ".wma", ".midi"],
                "Arşivler": [".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz", ".iso"],
                "Uygulamalar": [".exe", ".msi", ".dmg", ".apk", ".deb", ".rpm", ".appimage"],
                "Yazılım_Kod": [".py", ".html", ".css", ".js", ".cpp", ".c", ".java", ".php", ".rb", ".go", ".rs", ".swift"],
                "Veritabanı": [".db", ".sql", ".sqlite", ".mdb", ".accdb"],
                "E-Kitaplar": [".epub", ".mobi", ".azw", ".azw3"],
                "Fontlar": [".ttf", ".otf", ".woff", ".woff2"],
                "3D_Modeller": [".stl", ".obj", ".fbx", ".blend", ".3ds", ".dae"],
                "CAD_Dosyaları": [".dwg", ".dxf", ".skp"],
                "Tasarım": [".psd", ".ai", ".indd", ".xd", ".fig", ".sketch"],
                "Sistem": [".dll", ".sys", ".ini", ".cfg", ".bat", ".sh", ".reg"]
            },
            "document_categories": {
                "Finans": [".pdf", ".docx", ".xlsx", ".pptx"],
                "Eğitim": [".pdf", ".docx", ".pptx", ".xlsx"],
                "İş": [".docx", ".pdf", ".pptx", ".xlsx"],
                "Teknik": [".pdf", ".docx", ".txt", ".md"],
                "Sağlık": [".pdf", ".docx", ".xlsx"],
                "Hukuk": [".pdf", ".docx"],
                "Kişisel": [".docx", ".txt", ".pdf"],
                "Araştırma": [".pdf", ".docx"],
                "Tasarım": [".pdf", ".pptx", ".docx"],
                "Yönetim": [".pdf", ".docx", ".pptx", ".xlsx"]
            },
            "max_file_size": 500 * 1024 * 1024,
            "ignore_patterns": [".git", ".svn", ".idea", "__pycache__", "Thumbs.db", ".DS_Store", "desktop.ini"],
            "theme": "light",
            "language": "tr",
            "ai_enabled": True,
            "document_ai_enabled": True,
            "backup_enabled": True,
            "thread_count": 8,
            "ai_model": "Salesforce/blip-image-captioning-base",
            "multi_object_detection": True,
            "object_threshold": 30,
            "max_objects": 3,
            "enable_system_optimization": True,
            "document_analysis_depth": "medium"  # low, medium, high
        }
    
    def load_config(self) -> dict:
        """Yapılandırmayı yükle"""
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    loaded = json.load(f)
                    config = self.default_config.copy()
                    config.update(loaded)
                    return config
        except Exception as e:
            logging.error(f"Config yükleme hatası: {e}")
        return self.default_config.copy()
    
    def save_config(self, config: dict):
        """Yapılandırmayı kaydet"""
        try:
            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(config, f, indent=4, ensure_ascii=False)
        except Exception as e:
            logging.error(f"Config kaydetme hatası: {e}")

class Translator:
    """Çoklu dil desteği için sınıf"""
    
    def __init__(self):
        self.translations = {
            "tr": {
                "app_title": "AI Doküman Düzenleyici",
                "main": "Ana",
                "select_folder": "📁 Klasör Seç",
                "view_folder": "👁️ Klasörü Görüntüle",
                "start": "🚀 Düzenlemeyi Başlat",
                "ai_mode": "🤖 AI ile Düzenle",
                "document_ai_mode": "📄 AI ile Doküman Analizi",
                "undo": "↩️ Geri Al",
                "stats": "📊 İstatistikler",
                "settings": "⚙️ Ayarlar",
                "documents": "📄 Dokümanlar",
                "processing": "İşleniyor...",
                "completed": "Tamamlandı",
                "files": "Dosyalar",
                "folders": "Klasörler",
                "size": "Boyut",
                "time": "Süre",
                "details": "Detaylar",
                "about": "Hakkında",
                "folder_selection": "Klasör Seçimi",
                "options": "Seçenekler",
                "scan_subfolders": "Alt Klasörleri Tara",
                "delete_empty": "Boş Klasörleri Sil",
                "create_backup": "Yedek Oluştur",
                "language": "Dil",
                "theme": "Tema",
                "thread_count": "Thread Sayısı",
                "ai_enabled": "AI Özelliğini Aktif Et",
                "document_ai_enabled": "Doküman AI Analizi",
                "save_settings": "Ayarları Kaydet",
                "multi_object": "Çoklu Nesne Algılama",
                "object_threshold": "Nesne Eşik Değeri (%)",
                "max_objects": "Maksimum Nesne Sayısı",
                "ai_results": "AI Analiz Sonuçları",
                "file_history": "Dosya İşlem Geçmişi",
                "document_analysis": "Doküman Analizi",
                "document_summary": "Doküman Özeti",
                "document_categories": "Doküman Kategorileri",
                "general_settings": "Genel Ayarlar",
                "ai_settings": "AI Ayarları",
                "document_settings": "Doküman Ayarları",
                "thread_info": "İşlemcinize göre uygun thread sayısını seçin (i3: 2-4, i5/i7: 4-8, i9/Ryzen: 8-16)",
                "system_optimization": "Sistem Optimizasyonu",
                "system_optimization_tooltip": "Düşük donanımlı sistemlerde thread sayısını otomatik optimize eder",
                "analyze_documents": "📊 Dokümanları Analiz Et",
                "document_stats": "Doküman İstatistikleri",
                "content_analysis": "İçerik Analizi",
                "extract_text": "Metin Çıkar",
                "category_distribution": "Kategori Dağılımı",
                "keyword_analysis": "Anahtar Kelime Analizi",
                "view_document": "👁️ Dokümanı Görüntüle",
                "export_analysis": "📤 Analizi Dışa Aktar"
            },
            "en": {
                "app_title": "AI Document Organizer",
                "main": "Main",
                "select_folder": "📁 Select Folder",
                "view_folder": "👁️ View Folder",
                "start": "🚀 Start Organization",
                "ai_mode": "🤖 Organize with AI",
                "document_ai_mode": "📄 AI Document Analysis",
                "undo": "↩️ Undo",
                "stats": "📊 Statistics",
                "settings": "⚙️ Settings",
                "documents": "📄 Documents",
                "processing": "Processing...",
                "completed": "Completed",
                "files": "Files",
                "folders": "Folders",
                "size": "Size",
                "time": "Time",
                "details": "Details",
                "about": "About",
                "folder_selection": "Folder Selection",
                "options": "Options",
                "scan_subfolders": "Scan Subfolders",
                "delete_empty": "Delete Empty Folders",
                "create_backup": "Create Backup",
                "language": "Language",
                "theme": "Theme",
                "thread_count": "Thread Count",
                "ai_enabled": "Enable AI Feature",
                "document_ai_enabled": "Document AI Analysis",
                "save_settings": "Save Settings",
                "multi_object": "Multi-Object Detection",
                "object_threshold": "Object Threshold (%)",
                "max_objects": "Maximum Objects",
                "ai_results": "AI Analysis Results",
                "file_history": "File Processing History",
                "document_analysis": "Document Analysis",
                "document_summary": "Document Summary",
                "document_categories": "Document Categories",
                "general_settings": "General Settings",
                "ai_settings": "AI Settings",
                "document_settings": "Document Settings",
                "thread_info": "Select appropriate thread count for your CPU (i3: 2-4, i5/i7: 4-8, i9/Ryzen: 8-16)",
                "system_optimization": "System Optimization",
                "system_optimization_tooltip": "Automatically optimizes thread count on low-end systems",
                "analyze_documents": "📊 Analyze Documents",
                "document_stats": "Document Statistics",
                "content_analysis": "Content Analysis",
                "extract_text": "Extract Text",
                "category_distribution": "Category Distribution",
                "keyword_analysis": "Keyword Analysis",
                "view_document": "👁️ View Document",
                "export_analysis": "📤 Export Analysis"
            }
        }
    
    def get_text(self, key: str, lang: str = "tr") -> str:
        """Anahtar kelimenin çevirisini al"""
        return self.translations.get(lang, {}).get(key, key)

# ==================== SIGNAL CLASSES ====================

class WorkerSignals(QObject):
    """Çalışan thread'ler için sinyaller"""
    progress = Signal(int)
    status = Signal(str)
    finished = Signal(bool, dict)
    error = Signal(str)
    stats_update = Signal(dict)
    file_processed = Signal(str, str)
    ai_analysis_result = Signal(str, str, list)
    system_info = Signal(str)
    document_analyzed = Signal(str, str, str, str)  # filename, category, summary, metadata

# ==================== WORKER THREAD ====================

class OrganizerWorker(QThread):
    """Dosya düzenleme işçi thread'i"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.signals = WorkerSignals()
        self.config = {}
        self.ai_modu = False
        self.document_ai_modu = False
        self.include_subfolders = False
        self.delete_empty = False
        self.source_folder = ""
        self.ai_pipeline = None
        self.document_analyzer = None
        self.start_time = 0
        self.thread_count = 1
        self.enable_optimization = True
        self.stats = {
            "total_files": 0,
            "processed_files": 0,
            "skipped_files": 0,
            "created_folders": 0,
            "total_size": 0,
            "category_distribution": {},
            "detected_objects": {},
            "document_categories": {},
            "errors": []
        }
        self.log_entries = []
        self.processed_counter = 0
        self.lock = multiprocessing.Lock()
    
    def setup(self, config: dict, ai_modu: bool, document_ai_modu: bool, include_subfolders: bool, 
              delete_empty: bool, source_folder: str, thread_count: int = 1):
        """Worker'ı ayarla"""
        self.config = config
        self.ai_modu = ai_modu
        self.document_ai_modu = document_ai_modu
        self.include_subfolders = include_subfolders
        self.delete_empty = delete_empty
        self.source_folder = source_folder
        self.thread_count = thread_count
        self.enable_optimization = config.get("enable_system_optimization", True)
        
        # Doküman analizini başlat
        if self.document_ai_modu:
            self.document_analyzer = DocumentAnalyzer(config)
        
        self.stats = {
            "total_files": 0,
            "processed_files": 0,
            "skipped_files": 0,
            "created_folders": 0,
            "total_size": 0,
            "category_distribution": {},
            "detected_objects": {},
            "document_categories": {},
            "errors": []
        }
        self.log_entries = []
        self.processed_counter = 0
    
    def run(self):
        """Ana işlem döngüsü"""
        try:
            self.start_time = time.time()
            
            if self.ai_modu and self.config.get("ai_enabled", True):
                self._load_ai_pipeline()
            
            files = self._collect_files()
            self.stats["total_files"] = len(files)
            
            if not files:
                self.signals.status.emit("Düzenlenecek dosya bulunamadı!")
                self.signals.finished.emit(False, self.stats)
                return
            
            # Sistem bilgilerini göster
            self._show_system_info()
            
            self._process_files_multithreaded(files)
            
            if self.delete_empty:
                self._clean_empty_folders()
            
            self._save_log()
            self._calculate_statistics()
            
            self.signals.finished.emit(True, self.stats)
            
        except Exception as e:
            self.signals.error.emit(f"Kritik hata: {str(e)}")
            self.signals.finished.emit(False, self.stats)
    
    def _show_system_info(self):
        """Sistem bilgilerini göster"""
        if not PSUTIL_AVAILABLE:
            self.signals.system_info.emit("Sistem optimizasyonu devre dışı (psutil kurulu değil)")
            return
        
        try:
            cpu_count = psutil.cpu_count(logical=True) or 2
            cpu_physical = psutil.cpu_count(logical=False) or cpu_count // 2
            memory = psutil.virtual_memory()
            memory_gb = memory.total / (1024**3)
            
            info = (f"Sistem: {cpu_physical} çekirdek / {cpu_count} thread, "
                   f"RAM: {memory_gb:.1f}GB, "
                   f"İstenen thread: {self.thread_count}")
            
            self.signals.system_info.emit(info)
            logging.info(f"Sistem bilgisi: {info}")
            
        except Exception as e:
            logging.error(f"Sistem bilgisi alma hatası: {e}")
    
    def _optimize_thread_count(self, total_files: int, requested_threads: int) -> int:
        """Thread sayısını optimize et"""
        if not self.enable_optimization or not PSUTIL_AVAILABLE:
            # Optimizasyon kapalıysa veya psutil yoksa
            return min(requested_threads, 32)
        
        try:
            # Sistem bilgilerini al
            cpu_count = psutil.cpu_count(logical=True) or 2
            cpu_physical = psutil.cpu_count(logical=False) or max(1, cpu_count // 2)
            
            # CPU kullanımı
            cpu_percent = psutil.cpu_percent(interval=0.1)
            
            # Bellek durumu
            memory = psutil.virtual_memory()
            memory_percent = memory.percent
            memory_gb = memory.total / (1024**3)
            
            # İşlemci tipine göre sınıflandırma
            cpu_type = "unknown"
            if cpu_count <= 4:
                cpu_type = "low_end"  # i3, Celeron, Pentium
            elif cpu_count <= 8:
                cpu_type = "mid_range"  # i5, i7
            else:
                cpu_type = "high_end"  # i9, Ryzen
            
            # Başlangıç değeri
            optimal_threads = requested_threads
            
            # KURAL 1: Dosya sayısına göre sınırla
            if total_files < optimal_threads * 2:  # Her thread'e en az 2 dosya
                optimal_threads = max(1, total_files // 2)
                logging.info(f"Kural 1: Dosya sayısı az ({total_files}), thread {optimal_threads}")
            
            # KURAL 2: İşlemci tipine göre maksimum sınır
            max_by_cpu_type = {
                "low_end": min(4, cpu_count * 2),  # i3 için max 4
                "mid_range": min(8, cpu_count * 2),  # i5/i7 için max 8
                "high_end": min(16, cpu_count * 2),  # i9/Ryzen için max 16
                "unknown": min(8, cpu_count * 2)
            }
            
            cpu_limit = max_by_cpu_type.get(cpu_type, 8)
            if optimal_threads > cpu_limit:
                logging.info(f"Kural 2: {cpu_type} işlemci, thread {optimal_threads} -> {cpu_limit}")
                optimal_threads = cpu_limit
            
            # KURAL 3: CPU kullanımı yüksekse azalt
            if cpu_percent > 70:  # CPU %70'ten fazla kullanılıyorsa
                reduction = int(optimal_threads * 0.6)  # %40 azalt
                optimal_threads = max(2, reduction)
                logging.info(f"Kural 3: CPU %{cpu_percent}, thread {optimal_threads}")
            
            # KURAL 4: Bellek yüksekse azalt
            if memory_percent > 80:  # Bellek %80'den fazla kullanılıyorsa
                reduction = int(optimal_threads * 0.7)  # %30 azalt
                optimal_threads = max(2, reduction)
                logging.info(f"Kural 4: Memory %{memory_percent}, thread {optimal_threads}")
            
            # KURAL 5: Düşük RAM (<4GB) ise sınırla
            if memory_gb < 4:
                optimal_threads = min(optimal_threads, 4)
                logging.info(f"Kural 5: RAM {memory_gb:.1f}GB, thread max 4")
            
            # KURAL 6: Minimum ve maksimum sınırlar
            optimal_threads = max(1, min(optimal_threads, 32))
            
            # Log kaydı
            logging.info(f"Thread optimizasyonu: "
                        f"İstenen={requested_threads}, "
                        f"Optimal={optimal_threads}, "
                        f"CPU={cpu_type}({cpu_physical}/{cpu_count}), "
                        f"CPU%={cpu_percent}, "
                        f"Memory={memory_gb:.1f}GB({memory_percent}%), "
                        f"Files={total_files}")
            
            return optimal_threads
            
        except Exception as e:
            logging.error(f"Thread optimizasyon hatası: {e}")
            # Hata durumunda güvenli değer
            return min(requested_threads, 4)
    
    def _load_ai_pipeline(self):
        """AI pipeline'ını yükle"""
        try:
            self.signals.status.emit("AI Model yükleniyor...")
            self.ai_pipeline = pipeline("image-to-text", 
                                      model=self.config.get("ai_model", "Salesforce/blip-image-captioning-base"))
        except Exception as e:
            self.signals.error.emit(f"AI Model yüklenemedi: {str(e)}")
            self.ai_pipeline = None
    
    def _collect_files(self) -> List[str]:
        """Dosyaları topla"""
        files = []
        ignore_patterns = set(self.config.get("ignore_patterns", []))
        max_size = self.config.get("max_file_size", 500 * 1024 * 1024)
        
        try:
            if self.include_subfolders:
                for root, dirs, filenames in os.walk(self.source_folder):
                    dirs[:] = [d for d in dirs if d not in ignore_patterns]
                    
                    for filename in filenames:
                        if filename in ignore_patterns:
                            continue
                        
                        filepath = os.path.join(root, filename)
                        try:
                            if os.path.getsize(filepath) > max_size:
                                self.stats["skipped_files"] += 1
                                continue
                            files.append(filepath)
                        except OSError:
                            self.stats["skipped_files"] += 1
            else:
                for item in os.listdir(self.source_folder):
                    if item in ignore_patterns:
                        continue
                    
                    itempath = os.path.join(self.source_folder, item)
                    if os.path.isfile(itempath):
                        try:
                            if os.path.getsize(itempath) <= max_size:
                                files.append(itempath)
                            else:
                                self.stats["skipped_files"] += 1
                        except OSError:
                            self.stats["skipped_files"] += 1
            
            return files
        except Exception as e:
            self.signals.error.emit(f"Dosya toplama hatası: {str(e)}")
            return []
    
    def _process_files_multithreaded(self, files: List[str]):
        """Dosyaları çoklu thread ile işle"""
        total = len(files)
        if total == 0:
            return
        
        # Thread sayısını optimize et
        optimal_threads = self._optimize_thread_count(total, self.thread_count)
        
        # Eğer çok düşükse kullanıcıyı bilgilendir
        if optimal_threads < self.thread_count and self.enable_optimization:
            reduction_percent = int((1 - (optimal_threads / self.thread_count)) * 100)
            if PSUTIL_AVAILABLE:
                try:
                    cpu_count = psutil.cpu_count(logical=True) or 2
                    warning_msg = (f"Sistem optimizasyonu: "
                                  f"{self.thread_count} thread → {optimal_threads} thread "
                                  f"(CPU: {cpu_count} thread, %{reduction_percent} azaltma)")
                except:
                    warning_msg = f"Sistem optimizasyonu: Thread sayısı {optimal_threads}'a düşürüldü"
            else:
                warning_msg = f"Dosya sayısı az: Thread sayısı {optimal_threads}'a düşürüldü"
            
            self.signals.status.emit(warning_msg)
            logging.warning(warning_msg)
        
        thread_count = optimal_threads
        
        # Thread bilgisini gönder
        thread_info = f" ({thread_count} thread kullanılıyor)"
        self.signals.status.emit(f"Dosyalar işleniyor{thread_info}...")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=thread_count) as executor:
            # Tüm dosyaları thread'lere dağıt
            futures = []
            for i in range(thread_count):
                # Dosyaları thread'lere böl
                chunk_size = (total + thread_count - 1) // thread_count
                start_idx = i * chunk_size
                end_idx = min(start_idx + chunk_size, total)
                
                if start_idx < end_idx:
                    file_chunk = files[start_idx:end_idx]
                    future = executor.submit(self._process_file_chunk, file_chunk, total, start_idx)
                    futures.append(future)
            
            # Tüm thread'lerin bitmesini bekle
            for future in concurrent.futures.as_completed(futures):
                try:
                    result = future.result()
                except Exception as e:
                    self.signals.error.emit(f"Thread hatası: {str(e)}")
        
        # Tüm işlemler bitti, istatistikleri güncelle
        self.signals.stats_update.emit(self.stats)
    
    def _process_file_chunk(self, file_chunk: List[str], total_files: int, start_index: int):
        """Dosya grubunu işle"""
        for i, filepath in enumerate(file_chunk):
            try:
                filename = os.path.basename(filepath)
                ext = os.path.splitext(filename)[1].lower()
                
                # İlerleme durumunu güncelle
                with self.lock:
                    self.processed_counter += 1
                    # Her dosyada progress güncelleme
                    progress = int((self.processed_counter / total_files) * 100)
                    # Progress sinyalini lock dışında gönder
                    self.signals.progress.emit(progress)
                    
                    # Her 10 dosyada bir durum güncelle
                    if self.processed_counter % 10 == 0:
                        self.signals.status.emit(f"İşleniyor: {self.processed_counter}/{total_files} dosya")
                
                # Doküman analizi modu
                if self.document_ai_modu and ext in ['.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.txt', '.md']:
                    category, subcategory, detected_objects = self._determine_document_category(filepath)
                    
                    if category:
                        # Doküman istatistiklerini güncelle
                        with self.lock:
                            self.stats["document_categories"][category] = \
                                self.stats["document_categories"].get(category, 0) + 1
                    
                    target_category = category if category else "Diğer_Belgeler"
                    
                else:
                    # Normal AI modu
                    category, subcategory, detected_objects = self._determine_category(filepath)
                    target_category = category
                
                if subcategory:
                    target_dir = os.path.join(self.source_folder, target_category, subcategory)
                else:
                    target_dir = os.path.join(self.source_folder, target_category)
                
                os.makedirs(target_dir, exist_ok=True)
                
                new_path = self._move_file(filepath, target_dir, filename)
                
                if new_path:
                    with self.lock:
                        self.stats["processed_files"] += 1
                        file_size = os.path.getsize(new_path)
                        self.stats["total_size"] += file_size
                        
                        full_category = f"{target_category}/{subcategory}" if subcategory else target_category
                        self.stats["category_distribution"][full_category] = \
                            self.stats["category_distribution"].get(full_category, 0) + 1
                        
                        for obj in detected_objects:
                            self.stats["detected_objects"][obj] = \
                                self.stats["detected_objects"].get(obj, 0) + 1
                        
                        self.log_entries.append({
                            "old_path": filepath,
                            "new_path": new_path,
                            "category": full_category,
                            "objects": detected_objects,
                            "timestamp": datetime.now().isoformat()
                        })
                    
                    self.signals.file_processed.emit(filename, full_category)
                    
                    if detected_objects:
                        self.signals.ai_analysis_result.emit(filename, full_category, detected_objects)
                    
            except Exception as e:
                error_msg = f"{filename}: {str(e)}"
                with self.lock:
                    self.stats["errors"].append(error_msg)
                    self.stats["skipped_files"] += 1
    
    def _determine_category(self, filepath: str) -> Tuple[str, str, List[str]]:
        """Dosya için kategori belirle"""
        filename = os.path.basename(filepath)
        ext = os.path.splitext(filename)[1].lower()
        
        categories = self.config.get("kategoriler", {})
        for category, extensions in categories.items():
            if ext in extensions:
                if self.ai_modu and ext in ['.jpg', '.jpeg', '.png', '.webp', '.bmp']:
                    ai_category, detected_objects = self._ai_analyze(filepath)
                    return category, ai_category, detected_objects
                return category, "", []
        
        if self.ai_modu and self.ai_pipeline and ext in ['.jpg', '.jpeg', '.png', '.webp', '.bmp']:
            ai_category, detected_objects = self._ai_analyze(filepath)
            return "AI_Sınıflandırma", ai_category, detected_objects
        
        return "Diğerleri", "", []
    
    def _determine_document_category(self, filepath: str) -> Tuple[str, str, List[str]]:
        """Doküman için kategori belirle"""
        if not self.document_analyzer:
            return "Belgeler", "", []
        
        try:
            # Metni çıkar
            text = self.document_analyzer.extract_text(filepath)
            if not text:
                logging.info(f"Doküman metin çıkarılamadı: {filepath}")
                return "Belgeler", "", []
            
            logging.info(f"Doküman analiz ediliyor: {os.path.basename(filepath)} - {len(text)} karakter")
            
            # İçeriği analiz et
            category, related_categories, summary = self.document_analyzer.analyze_content(text)
            
            # Metadata çıkar
            metadata = self.document_analyzer.extract_metadata(filepath, text)
            
            # Sinyal gönder
            self.signals.document_analyzed.emit(
                os.path.basename(filepath),
                category,
                summary,
                json.dumps(metadata, ensure_ascii=False, indent=2)
            )
            
            # İlgili kategorileri alt kategori olarak kullan
            subcategory = ""
            if related_categories:
                subcategory = "_".join(related_categories[:2])
            
            return category, subcategory, []
            
        except Exception as e:
            logging.error(f"Doküman analiz hatası {filepath}: {e}")
            return "Belgeler", "", []
    
    def _ai_analyze(self, filepath: str) -> Tuple[str, List[str]]:
        """AI ile görsel analizi"""
        try:
            if self.ai_pipeline:
                result = self.ai_pipeline(filepath)
                caption = result[0]['generated_text'].lower()
                
                detected_objects = self._detect_objects(caption)
                category = self._determine_ai_category(detected_objects, caption)
                
                return category, detected_objects
        except Exception as e:
            logging.error(f"AI analiz hatası {filepath}: {e}")
        
        return "Bilinmeyen", []
    
    def _detect_objects(self, caption: str) -> List[str]:
        """Caption'dan nesneleri tespit et"""
        objects = []
        
        # Eşik değerini al (% olarak 1-100)
        threshold_percent = self.config.get("object_threshold", 30)
        threshold = threshold_percent / 100.0  # 0.0-1.0 arası
        
        # Genişletilmiş nesne tanıma anahtar kelimeleri
        object_keywords = [
            "car", "truck", "vehicle", "automobile", "bus", "motorcycle", "bicycle", "bike",
            "chair", "table", "desk", "furniture", "sofa", "couch", "bed",
            "computer", "laptop", "phone", "smartphone", "tablet", "monitor", "screen",
            "building", "house", "home", "apartment", "skyscraper", "office",
            "tree", "plant", "flower", "forest", "wood", "nature",
            "mountain", "hill", "valley", "landscape", "scenery",
            "food", "fruit", "vegetable", "meal", "dish", "drink", "water", "coffee", "tea",
            "clothing", "shirt", "pants", "dress", "shoe", "jacket", "hat",
            "sport", "football", "soccer", "basketball", "tennis", "golf", "swimming",
            "music", "guitar", "piano", "violin", "drum", "instrument", "song",
            "book", "document", "paper", "text", "writing", "letter",
            "tool", "hammer", "screwdriver", "wrench", "drill", "saw",
            "art", "painting", "drawing", "sculpture", "statue", "picture", "photo",
            "person", "man", "woman", "child", "baby", "people", "human",
            "animal", "dog", "cat", "bird", "fish", "horse", "elephant", "lion",
            "sky", "cloud", "sun", "moon", "star", "weather", "rain",
            "water", "sea", "ocean", "river", "lake", "beach", "wave",
            "city", "street", "road", "highway", "bridge", "park", "garden"
        ]
        
        caption_lower = caption.lower()
        caption_words = caption_lower.split()
        
        for keyword in object_keywords:
            # Anahtar kelimenin caption'da bulunma oranını hesapla
            keyword_words = keyword.split()
            match_score = 0
            
            for kw_word in keyword_words:
                if len(kw_word) > 3:  # Kısa kelimeleri atla
                    # Tam eşleşme kontrolü
                    if kw_word in caption_words:
                        match_score += 1.0
                    # Kısmi eşleşme kontrolü (kelimenin bir kısmı)
                    elif any(kw_word in word for word in caption_words if len(word) > 3):
                        match_score += 0.5
            
            # Eşik değeri kontrolü
            if keyword_words:
                match_ratio = match_score / len(keyword_words)
                if match_ratio >= threshold:
                    objects.append(keyword)
        
        max_objects = self.config.get("max_objects", 3)
        return objects[:max_objects]
    
    def _determine_ai_category(self, objects: List[str], caption: str) -> str:
        """Nesnelere göre AI kategorisi belirle"""
        if not objects:
            # Nesne yoksa caption'dan kelimeler al
            words = [word for word in caption.split() if len(word) > 3]
            if words:
                return "_".join(words[:2]).capitalize()[:20]
            return "Genel"
        
        if len(objects) > 1:
            return f"{objects[0]}_{objects[1]}".capitalize()
        
        return objects[0].capitalize()
    
    def _move_file(self, old_path: str, target_dir: str, filename: str) -> Optional[str]:
        """Dosyayı taşı"""
        new_path = os.path.join(target_dir, filename)
        
        counter = 1
        name, ext = os.path.splitext(filename)
        while os.path.exists(new_path):
            new_path = os.path.join(target_dir, f"{name}_{counter}{ext}")
            counter += 1
        
        try:
            shutil.move(old_path, new_path)
            return new_path
        except Exception as e:
            raise Exception(f"Taşıma hatası: {str(e)}")
    
    def _clean_empty_folders(self):
        """Boş klasörleri temizle"""
        try:
            for root, dirs, files in os.walk(self.source_folder, topdown=False):
                for dir_name in dirs:
                    dir_path = os.path.join(root, dir_name)
                    try:
                        if not os.listdir(dir_path):
                            os.rmdir(dir_path)
                    except OSError:
                        pass
        except Exception as e:
            logging.error(f"Boş klasör temizleme hatası: {e}")
    
    def _save_log(self):
        """İşlem log'unu kaydet"""
        if not self.log_entries:
            return
        
        log_path = os.path.join(self.source_folder, ".dosya_sirala_duzenle_gecmis.json")
        
        try:
            existing_logs = []
            if os.path.exists(log_path):
                with open(log_path, 'r', encoding='utf-8') as f:
                    existing_logs = json.load(f)
            
            new_log_entry = {
                "timestamp": datetime.now().isoformat(),
                "ai_mode": self.ai_modu,
                "document_ai_mode": self.document_ai_modu,
                "requested_threads": self.thread_count,
                "actual_threads": self.stats.get("thread_count_used", 1),
                "optimization_enabled": self.enable_optimization,
                "total_files": self.stats["total_files"],
                "processed_files": self.stats["processed_files"],
                "document_categories": self.stats.get("document_categories", {}),
                "processing_time": self.stats.get("processing_time", 0),
                "movements": self.log_entries
            }
            
            existing_logs.append(new_log_entry)
            
            if len(existing_logs) > 5:
                existing_logs = existing_logs[-5:]
            
            with open(log_path, 'w', encoding='utf-8') as f:
                json.dump(existing_logs, f, indent=4, ensure_ascii=False)
                
        except Exception as e:
            logging.error(f"Log kaydetme hatası: {e}")
    
    def _calculate_statistics(self):
        """İstatistikleri hesapla"""
        self.stats["processing_time"] = time.time() - self.start_time
        self.stats["thread_count_used"] = self.thread_count
        
        folder_count = 0
        try:
            for root, dirs, files in os.walk(self.source_folder):
                folder_count += len(dirs)
        except:
            pass
        
        self.stats["created_folders"] = folder_count
        
        size = self.stats["total_size"]
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size < 1024.0:
                self.stats["formatted_size"] = f"{size:.2f} {unit}"
                break
            size /= 1024.0

# ==================== DOCUMENT ANALYSIS THREAD ====================

class DocumentAnalysisWorker(QThread):
    """Doküman analizi için özel thread"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.signals = WorkerSignals()
        self.files = []
        self.config = {}
        self.analyzer = None
        self.results = []
    
    def setup(self, files: List[str], config: dict):
        """Worker'ı ayarla"""
        self.files = files
        self.config = config
        self.analyzer = DocumentAnalyzer(config)
        self.results = []
    
    def run(self):
        """Doküman analizini çalıştır"""
        try:
            total = len(self.files)
            for i, filepath in enumerate(self.files):
                try:
                    filename = os.path.basename(filepath)
                    
                    # İlerleme durumunu güncelle
                    progress = int((i + 1) / total * 100)
                    self.signals.progress.emit(progress)
                    
                    # Metni çıkar
                    text = self.analyzer.extract_text(filepath)
                    if not text:
                        continue
                    
                    # İçeriği analiz et
                    category, related_categories, summary = self.analyzer.analyze_content(text, filename)
                    
                    # Metadata çıkar
                    metadata = self.analyzer.extract_metadata(filepath, text)
                    
                    # Sonuçları kaydet
                    result = {
                        "filename": filename,
                        "path": filepath,
                        "category": category,
                        "related_categories": related_categories,
                        "summary": summary,
                        "metadata": metadata,
                        "text_preview": text[:500] + "..." if len(text) > 500 else text
                    }
                    self.results.append(result)
                    
                    # Sinyal gönder
                    self.signals.document_analyzed.emit(
                        filename,
                        category,
                        summary,
                        json.dumps(metadata, ensure_ascii=False)
                    )
                    
                    # Durum güncelle
                    if (i + 1) % 5 == 0:
                        self.signals.status.emit(f"Analiz ediliyor: {i + 1}/{total} doküman")
                    
                except Exception as e:
                    logging.error(f"Doküman analiz hatası {filepath}: {e}")
            
            # Tamamlandı sinyali
            self.signals.finished.emit(True, {"total_analyzed": len(self.results), "results": self.results})
            
        except Exception as e:
            self.signals.error.emit(f"Doküman analiz hatası: {str(e)}")
            self.signals.finished.emit(False, {})

# ==================== MAIN WINDOW ====================

class DosyaDuzenleyici(QMainWindow):
    """Ana uygulama penceresi"""
    
    def __init__(self):
        super().__init__()
        
        self.setup_logging()
        self.config_manager = ConfigManager()
        self.config = self.config_manager.load_config()
        self.translator = Translator()
        
        self.selected_folder = ""
        self.log_file = ".dosya_sirala_duzenle_gecmis.json"
        self.worker_thread = None
        self.document_worker = None
        self.current_language = self.config.get("language", "tr")
        self.current_theme = self.config.get("theme", "light")
        self.document_results = []
        
        self.init_ui()
        self.apply_theme()
        self.update_ui_texts()
    
    def setup_logging(self):
        """Loglama sistemini kur"""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler('dosya_duzenleyici.log', encoding='utf-8'),
                logging.StreamHandler()
            ]
        )
        self.logger = logging.getLogger(__name__)
    
    def init_ui(self):
        """Kullanıcı arayüzünü başlat"""
        self.setWindowTitle("AI Doküman Düzenleyici")
        self.setMinimumSize(1100, 800)
        
        # Merkezi widget
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)
        main_layout.setContentsMargins(10, 10, 10, 10)
        main_layout.setSpacing(10)
        
        # Sekmeli panel
        self.tab_widget = QTabWidget()
        main_layout.addWidget(self.tab_widget)
        
        self.setup_main_tab()
        self.setup_documents_tab()
        self.setup_stats_tab()
        self.setup_settings_tab()
        
        self.setup_status_bar()
    
    def setup_main_tab(self):
        """Ana sekme bileşenleri"""
        main_tab = QWidget()
        layout = QVBoxLayout(main_tab)
        layout.setContentsMargins(10, 10, 10, 10)
        layout.setSpacing(15)
        
        # Başlık
        self.title_label = QLabel("📁 AI Doküman Düzenleyici")
        self.title_label.setAlignment(Qt.AlignCenter)
        self.title_label.setProperty("cssClass", "title")
        layout.addWidget(self.title_label)
        
        # Klasör seçim alanı
        self.folder_group = QGroupBox("Klasör Seçimi")
        folder_layout = QVBoxLayout(self.folder_group)
        folder_layout.setSpacing(10)
        
        self.folder_label = QLabel("Henüz klasör seçilmedi")
        self.folder_label.setWordWrap(True)
        folder_layout.addWidget(self.folder_label)
        
        btn_layout = QHBoxLayout()
        btn_layout.setSpacing(10)
        
        self.btn_select = QPushButton("📁 Klasör Seç")
        self.btn_select.clicked.connect(self.select_folder)
        self.btn_select.setFixedHeight(40)
        
        self.btn_view = QPushButton("👁️ Klasörü Görüntüle")
        self.btn_view.clicked.connect(self.view_current_folder)
        self.btn_view.setEnabled(False)
        self.btn_view.setFixedHeight(40)
        
        btn_layout.addWidget(self.btn_select)
        btn_layout.addWidget(self.btn_view)
        btn_layout.addStretch()
        folder_layout.addLayout(btn_layout)
        
        layout.addWidget(self.folder_group)
        
        # Seçenekler
        self.options_group = QGroupBox("Seçenekler")
        options_layout = QVBoxLayout(self.options_group)
        options_layout.setSpacing(8)
        
        self.cb_subfolders = QCheckBox("Alt Klasörleri Tara")
        self.cb_subfolders.setChecked(True)
        
        self.cb_empty_folders = QCheckBox("Boş Klasörleri Sil")
        
        self.cb_backup = QCheckBox("Yedek Oluştur")
        self.cb_backup.setChecked(self.config.get("backup_enabled", True))
        
        options_layout.addWidget(self.cb_subfolders)
        options_layout.addWidget(self.cb_empty_folders)
        options_layout.addWidget(self.cb_backup)
        
        layout.addWidget(self.options_group)
        
        # Sistem bilgisi etiketi
        self.system_info_label = QLabel("")
        self.system_info_label.setAlignment(Qt.AlignCenter)
        self.system_info_label.setProperty("cssClass", "subtitle")
        self.system_info_label.setStyleSheet("color: #666; font-size: 11px;")
        self.system_info_label.setWordWrap(True)
        layout.addWidget(self.system_info_label)
        
        # İlerleme çubuğu
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        layout.addWidget(self.progress_bar)
        
        # Durum etiketi
        self.status_label = QLabel("")
        self.status_label.setAlignment(Qt.AlignCenter)
        self.status_label.setProperty("cssClass", "subtitle")
        layout.addWidget(self.status_label)
        
        # Kontrol butonları
        control_layout = QHBoxLayout()
        control_layout.setSpacing(10)
        
        self.btn_ai = QPushButton("🤖 AI ile Düzenle")
        self.btn_ai.clicked.connect(lambda: self.start_organization(ai_mode=True, document_ai=False))
        self.btn_ai.setEnabled(False)
        self.btn_ai.setProperty("cssClass", "info")
        
        self.btn_document_ai = QPushButton("📄 AI Doküman Analizi")
        self.btn_document_ai.clicked.connect(lambda: self.start_organization(ai_mode=False, document_ai=True))
        self.btn_document_ai.setEnabled(False)
        self.btn_document_ai.setProperty("cssClass", "warning")
        
        self.btn_start = QPushButton("🚀 Düzenlemeyi Başlat")
        self.btn_start.clicked.connect(lambda: self.start_organization(ai_mode=False, document_ai=False))
        self.btn_start.setEnabled(False)
        self.btn_start.setProperty("cssClass", "success")
        
        self.btn_undo = QPushButton("↩️ Geri Al")
        self.btn_undo.clicked.connect(self.undo_action)
        self.btn_undo.setEnabled(False)
        self.btn_undo.setProperty("cssClass", "danger")
        
        control_layout.addWidget(self.btn_ai)
        control_layout.addWidget(self.btn_document_ai)
        control_layout.addWidget(self.btn_start)
        control_layout.addWidget(self.btn_undo)
        
        layout.addLayout(control_layout)
        
        # AI Analiz Sonuçları
        self.ai_group = QGroupBox("AI Analiz Sonuçları")
        ai_layout = QVBoxLayout(self.ai_group)
        
        self.ai_results_text = QTextEdit()
        self.ai_results_text.setReadOnly(True)
        self.ai_results_text.setMaximumHeight(100)
        self.ai_results_text.setPlaceholderText("AI analiz sonuçları burada gösterilecek...")
        ai_layout.addWidget(self.ai_results_text)
        
        layout.addWidget(self.ai_group)
        
        # Dosya listesi
        self.file_group = QGroupBox("Dosya İşlem Geçmişi")
        file_layout = QVBoxLayout(self.file_group)
        
        self.file_list = QTextEdit()
        self.file_list.setReadOnly(True)
        self.file_list.setMaximumHeight(150)
        file_layout.addWidget(self.file_list)
        
        layout.addWidget(self.file_group)
        
        self.tab_widget.addTab(main_tab, "Ana")
    
    def setup_documents_tab(self):
        """Doküman analizi sekmesi"""
        documents_tab = QWidget()
        layout = QVBoxLayout(documents_tab)
        layout.setContentsMargins(10, 10, 10, 10)
        layout.setSpacing(15)
        
        # Başlık
        doc_title = QLabel("📄 Doküman Analizi ve Sınıflandırma")
        doc_title.setAlignment(Qt.AlignCenter)
        doc_title.setProperty("cssClass", "title")
        layout.addWidget(doc_title)
        
        # Kontrol paneli
        control_panel = QGroupBox("Analiz Kontrolleri")
        control_layout = QVBoxLayout(control_panel)
        
        # Analiz butonu
        self.btn_analyze_documents = QPushButton("📊 Dokümanları Analiz Et")
        self.btn_analyze_documents.clicked.connect(self.analyze_documents)
        self.btn_analyze_documents.setEnabled(False)
        self.btn_analyze_documents.setProperty("cssClass", "info")
        control_layout.addWidget(self.btn_analyze_documents)
        
        # Filtreleme
        filter_layout = QHBoxLayout()
        filter_layout.addWidget(QLabel("Dosya Türü:"))
        
        self.doc_filter_combo = QComboBox()
        self.doc_filter_combo.addItems(["Tümü", "PDF", "DOCX", "PPTX", "Excel", "TXT"])
        filter_layout.addWidget(self.doc_filter_combo)
        
        self.btn_filter = QPushButton("Filtrele")
        self.btn_filter.clicked.connect(self.filter_documents)
        self.btn_filter.setEnabled(False)
        filter_layout.addWidget(self.btn_filter)
        filter_layout.addStretch()
        
        control_layout.addLayout(filter_layout)
        layout.addWidget(control_panel)
        
        # Splitter ile iki bölme
        splitter = QSplitter(Qt.Vertical)
        
        # Üst bölme: Doküman listesi
        doc_list_widget = QWidget()
        doc_list_layout = QVBoxLayout(doc_list_widget)
        
        self.doc_tree = QTreeWidget()
        self.doc_tree.setHeaderLabels(["Dosya", "Kategori", "Özet", "Boyut"])
        self.doc_tree.setColumnWidth(0, 200)
        self.doc_tree.setColumnWidth(1, 150)
        self.doc_tree.setColumnWidth(2, 300)
        self.doc_tree.itemDoubleClicked.connect(self.view_document_details)
        doc_list_layout.addWidget(self.doc_tree)
        
        splitter.addWidget(doc_list_widget)
        
        # Alt bölme: Detaylar
        detail_widget = QWidget()
        detail_layout = QVBoxLayout(detail_widget)
        
        detail_tabs = QTabWidget()
        
        # Özet sekmesi
        summary_tab = QWidget()
        summary_layout = QVBoxLayout(summary_tab)
        self.doc_summary_text = QTextEdit()
        self.doc_summary_text.setReadOnly(True)
        self.doc_summary_text.setPlaceholderText("Doküman özeti burada gösterilecek...")
        summary_layout.addWidget(self.doc_summary_text)
        detail_tabs.addTab(summary_tab, "Özet")
        
        # Metadata sekmesi
        metadata_tab = QWidget()
        metadata_layout = QVBoxLayout(metadata_tab)
        self.doc_metadata_text = QTextEdit()
        self.doc_metadata_text.setReadOnly(True)
        self.doc_metadata_text.setPlaceholderText("Metadata bilgileri burada gösterilecek...")
        metadata_layout.addWidget(self.doc_metadata_text)
        detail_tabs.addTab(metadata_tab, "Metadata")
        
        # Metin önizleme
        preview_tab = QWidget()
        preview_layout = QVBoxLayout(preview_tab)
        self.doc_preview_text = QTextEdit()
        self.doc_preview_text.setReadOnly(True)
        self.doc_preview_text.setPlaceholderText("Metin önizlemesi burada gösterilecek...")
        preview_layout.addWidget(self.doc_preview_text)
        detail_tabs.addTab(preview_tab, "Metin Önizleme")
        
        detail_layout.addWidget(detail_tabs)
        
        # Alt butonlar
        button_layout = QHBoxLayout()
        self.btn_view_doc = QPushButton("👁️ Dokümanı Görüntüle")
        self.btn_view_doc.clicked.connect(self.open_document)
        self.btn_view_doc.setEnabled(False)
        
        self.btn_export_analysis = QPushButton("📤 Analizi Dışa Aktar")
        self.btn_export_analysis.clicked.connect(self.export_analysis)
        self.btn_export_analysis.setEnabled(False)
        
        button_layout.addWidget(self.btn_view_doc)
        button_layout.addWidget(self.btn_export_analysis)
        button_layout.addStretch()
        
        detail_layout.addLayout(button_layout)
        splitter.addWidget(detail_widget)
        
        # Splitter boyutları
        splitter.setSizes([400, 300])
        
        layout.addWidget(splitter)
        
        # İstatistikler
        stats_group = QGroupBox("Doküman İstatistikleri")
        stats_layout = QVBoxLayout(stats_group)
        
        self.doc_stats_text = QTextEdit()
        self.doc_stats_text.setReadOnly(True)
        self.doc_stats_text.setMaximumHeight(100)
        self.doc_stats_text.setPlaceholderText("Analiz istatistikleri burada gösterilecek...")
        stats_layout.addWidget(self.doc_stats_text)
        
        layout.addWidget(stats_group)
        
        self.tab_widget.addTab(documents_tab, "📄 Dokümanlar")
    
    def setup_stats_tab(self):
        """İstatistikler sekmesi"""
        stats_tab = QWidget()
        layout = QVBoxLayout(stats_tab)
        layout.setContentsMargins(10, 10, 10, 10)
        layout.setSpacing(15)
        
        # İstatistikler
        self.stats_group = QGroupBox("İstatistikler")
        stats_layout = QVBoxLayout(self.stats_group)
        
        self.stats_table = QTableWidget()
        self.stats_table.setColumnCount(2)
        self.stats_table.setHorizontalHeaderLabels(["Özellik", "Değer"])
        self.stats_table.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
        stats_layout.addWidget(self.stats_table)
        
        layout.addWidget(self.stats_group)
        
        # Kategori dağılımı
        self.category_group = QGroupBox("Kategori Dağılımı")
        category_layout = QVBoxLayout(self.category_group)
        
        self.category_table = QTableWidget()
        self.category_table.setColumnCount(2)
        self.category_table.setHorizontalHeaderLabels(["Kategori", "Dosya Sayısı"])
        self.category_table.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
        category_layout.addWidget(self.category_table)
        
        layout.addWidget(self.category_group)
        
        # Doküman kategorileri
        self.document_category_group = QGroupBox("Doküman Kategorileri")
        doc_category_layout = QVBoxLayout(self.document_category_group)
        
        self.document_category_table = QTableWidget()
        self.document_category_table.setColumnCount(2)
        self.document_category_table.setHorizontalHeaderLabels(["Kategori", "Doküman Sayısı"])
        self.document_category_table.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
        doc_category_layout.addWidget(self.document_category_table)
        
        layout.addWidget(self.document_category_group)
        
        # Nesneler
        self.objects_group = QGroupBox("Tespit Edilen Nesneler")
        objects_layout = QVBoxLayout(self.objects_group)
        
        self.objects_table = QTableWidget()
        self.objects_table.setColumnCount(2)
        self.objects_table.setHorizontalHeaderLabels(["Nesne", "Sayı"])
        self.objects_table.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
        objects_layout.addWidget(self.objects_table)
        
        layout.addWidget(self.objects_group)
        
        self.tab_widget.addTab(stats_tab, "📊 İstatistikler")
    
    def setup_settings_tab(self):
        """Ayarlar sekmesi"""
        settings_tab = QWidget()
        layout = QVBoxLayout(settings_tab)
        layout.setContentsMargins(10, 10, 10, 10)
        layout.setSpacing(15)
        
        # Genel Ayarlar
        self.general_group = QGroupBox("Genel Ayarlar")
        general_layout = QVBoxLayout(self.general_group)
        general_layout.setSpacing(10)
        
        # Dil seçimi
        lang_layout = QHBoxLayout()
        self.lang_label = QLabel("Dil:")
        self.lang_label.setMinimumWidth(120)
        self.cb_language = QComboBox()
        self.cb_language.addItems(["Türkçe", "English"])
        self.cb_language.setCurrentText("Türkçe" if self.current_language == "tr" else "English")
        self.cb_language.currentTextChanged.connect(self.change_language)
        lang_layout.addWidget(self.lang_label)
        lang_layout.addWidget(self.cb_language)
        lang_layout.addStretch()
        general_layout.addLayout(lang_layout)
        
        # Tema seçimi
        theme_layout = QHBoxLayout()
        self.theme_label = QLabel("Tema:")
        self.theme_label.setMinimumWidth(120)
        self.cb_theme = QComboBox()
        self.cb_theme.addItems(["Açık", "Koyu"])
        self.cb_theme.setCurrentText("Açık" if self.current_theme == "light" else "Koyu")
        self.cb_theme.currentTextChanged.connect(self.change_theme)
        theme_layout.addWidget(self.theme_label)
        theme_layout.addWidget(self.cb_theme)
        theme_layout.addStretch()
        general_layout.addLayout(theme_layout)
        
        # Thread sayısı
        thread_layout = QHBoxLayout()
        self.thread_label = QLabel("Thread Sayısı:")
        self.thread_label.setMinimumWidth(120)
        self.spin_threads = QSpinBox()
        self.spin_threads.setRange(1, 32)
        self.spin_threads.setValue(self.config.get("thread_count", 8))
        thread_layout.addWidget(self.thread_label)
        thread_layout.addWidget(self.spin_threads)
        thread_layout.addStretch()
        general_layout.addLayout(thread_layout)
        
        # Thread bilgisi etiketi
        self.thread_info_label = QLabel("İşlemcinize göre uygun thread sayısını seçin (i3: 2-4, i5/i7: 4-8, i9/Ryzen: 8-16)")
        self.thread_info_label.setStyleSheet("color: #666; font-size: 11px;")
        self.thread_info_label.setWordWrap(True)
        general_layout.addWidget(self.thread_info_label)
        
        # Sistem optimizasyonu
        self.cb_system_optimization = QCheckBox("Sistem Optimizasyonu")
        self.cb_system_optimization.setChecked(self.config.get("enable_system_optimization", True))
        self.cb_system_optimization.setToolTip("Düşük donanımlı sistemlerde thread sayısını otomatik optimize eder")
        general_layout.addWidget(self.cb_system_optimization)
        
        # psutil durumu
        if not PSUTIL_AVAILABLE:
            psutil_warning = QLabel("⚠️ psutil kütüphanesi kurulu değil. Sistem optimizasyonu devre dışı.")
            psutil_warning.setStyleSheet("color: #e74c3c; font-size: 11px; font-weight: bold;")
            psutil_warning.setWordWrap(True)
            general_layout.addWidget(psutil_warning)
        
        layout.addWidget(self.general_group)
        
        # AI Ayarları
        self.ai_settings_group = QGroupBox("AI Ayarları")
        ai_layout = QVBoxLayout(self.ai_settings_group)
        ai_layout.setSpacing(10)
        
        self.cb_ai_enabled = QCheckBox("AI Özelliğini Aktif Et")
        self.cb_ai_enabled.setChecked(self.config.get("ai_enabled", True))
        
        self.cb_document_ai_enabled = QCheckBox("Doküman AI Analizi")
        self.cb_document_ai_enabled.setChecked(self.config.get("document_ai_enabled", True))
        
        self.cb_multi_object = QCheckBox("Çoklu Nesne Algılama")
        self.cb_multi_object.setChecked(self.config.get("multi_object_detection", True))
        
        # Nesne Eşik Değeri - SLIDER ile 1-100 arası
        threshold_layout = QHBoxLayout()
        self.threshold_label = QLabel("Nesne Eşik Değeri (%):")
        self.threshold_label.setMinimumWidth(120)
        
        self.threshold_slider = QSlider(Qt.Horizontal)
        self.threshold_slider.setRange(1, 100)
        self.threshold_slider.setValue(self.config.get("object_threshold", 30))
        
        self.threshold_value_label = QLabel(f"%{self.config.get('object_threshold', 30)}")
        self.threshold_value_label.setMinimumWidth(40)
        
        self.threshold_slider.valueChanged.connect(
            lambda value: self.threshold_value_label.setText(f"%{value}")
        )
        
        threshold_layout.addWidget(self.threshold_label)
        threshold_layout.addWidget(self.threshold_slider)
        threshold_layout.addWidget(self.threshold_value_label)
        threshold_layout.addStretch()
        
        max_objects_layout = QHBoxLayout()
        self.max_objects_label = QLabel("Maksimum Nesne Sayısı:")
        self.max_objects_label.setMinimumWidth(120)
        self.spin_max_objects = QSpinBox()
        self.spin_max_objects.setRange(1, 10)
        self.spin_max_objects.setValue(self.config.get("max_objects", 3))
        max_objects_layout.addWidget(self.max_objects_label)
        max_objects_layout.addWidget(self.spin_max_objects)
        max_objects_layout.addStretch()
        
        ai_layout.addWidget(self.cb_ai_enabled)
        ai_layout.addWidget(self.cb_document_ai_enabled)
        ai_layout.addWidget(self.cb_multi_object)
        ai_layout.addLayout(threshold_layout)
        ai_layout.addLayout(max_objects_layout)
        
        layout.addWidget(self.ai_settings_group)
        
        # Doküman Ayarları
        self.document_settings_group = QGroupBox("Doküman Ayarları")
        doc_layout = QVBoxLayout(self.document_settings_group)
        doc_layout.setSpacing(10)
        
        # Analiz derinliği
        depth_layout = QHBoxLayout()
        depth_layout.addWidget(QLabel("Analiz Derinliği:"))
        
        self.analysis_depth_combo = QComboBox()
        self.analysis_depth_combo.addItems(["Düşük", "Orta", "Yüksek"])
        current_depth = self.config.get("document_analysis_depth", "medium")
        index = {"low": 0, "medium": 1, "high": 2}.get(current_depth, 1)
        self.analysis_depth_combo.setCurrentIndex(index)
        depth_layout.addWidget(self.analysis_depth_combo)
        depth_layout.addStretch()
        
        doc_layout.addLayout(depth_layout)
        
        # Desteklenen dosya türleri
        file_types_label = QLabel("Desteklenen Dosya Türleri: PDF, DOCX, PPTX, Excel, TXT")
        file_types_label.setStyleSheet("color: #666; font-size: 11px;")
        doc_layout.addWidget(file_types_label)
        
        # Kütüphane durumları
        lib_status = QLabel(self.get_library_status())
        lib_status.setStyleSheet("color: #666; font-size: 10px;")
        lib_status.setWordWrap(True)
        doc_layout.addWidget(lib_status)
        
        layout.addWidget(self.document_settings_group)
        
        # Kaydet butonu
        self.btn_save_settings = QPushButton("Ayarları Kaydet")
        self.btn_save_settings.clicked.connect(self.save_settings)
        self.btn_save_settings.setProperty("cssClass", "success")
        
        layout.addWidget(self.btn_save_settings)
        layout.addStretch()
        
        self.tab_widget.addTab(settings_tab, "⚙️ Ayarlar")
    
    def get_library_status(self):
        """Kütüphane durumlarını al"""
        status = []
        if PDF_AVAILABLE:
            status.append("✓ PDF")
        else:
            status.append("✗ PDF")
        
        if DOCX_AVAILABLE:
            status.append("✓ DOCX")
        else:
            status.append("✗ DOCX")
        
        if PPTX_AVAILABLE:
            status.append("✓ PPTX")
        else:
            status.append("✗ PPTX")
        
        if EXCEL_AVAILABLE:
            status.append("✓ Excel (.xlsx)")
        else:
            status.append("✗ Excel (.xlsx)")
        
        if XLRD_AVAILABLE:
            status.append("✓ Excel (.xls)")
        else:
            status.append("✗ Excel (.xls)")
        
        if TRANSFORMERS_AVAILABLE:
            status.append("✓ AI Model")
        else:
            status.append("✗ AI Model")
        
        return " | ".join(status)
    
    def setup_status_bar(self):
        """Durum çubuğunu ayarla"""
        self.statusBar().showMessage("Hazır")
    
    def apply_theme(self):
        """Temayı uygula"""
        theme = ThemeManager.get_theme(self.current_theme)
        stylesheet = ThemeManager.generate_stylesheet(theme)
        self.setStyleSheet(stylesheet)
    
    def select_folder(self):
        """Klasör seç"""
        folder_dialog_title = "Klasör Seç" if self.current_language == "tr" else "Select Folder"
        folder = QFileDialog.getExistingDirectory(self, folder_dialog_title)
        if folder:
            self.selected_folder = os.path.abspath(folder)
            folder_name_msg = f"Seçilen: {os.path.basename(self.selected_folder)}" if self.current_language == "tr" else f"Selected: {os.path.basename(self.selected_folder)}"
            self.folder_label.setText(folder_name_msg)
            self.btn_start.setEnabled(True)
            self.btn_ai.setEnabled(True)
            self.btn_document_ai.setEnabled(True)
            self.btn_analyze_documents.setEnabled(True)
            self.btn_view.setEnabled(True)
            self.check_log_file()
    
    def view_current_folder(self):
        """Mevcut klasörü görüntüle"""
        if self.selected_folder:
            QDesktopServices.openUrl(QUrl.fromLocalFile(self.selected_folder))
    
    def start_organization(self, ai_mode: bool, document_ai: bool):
        """Düzenleme işlemini başlat"""
        if not self.selected_folder:
            warning_msg = "Lütfen önce bir klasör seçin!" if self.current_language == "tr" else "Please select a folder first!"
            QMessageBox.warning(self, "Uyarı" if self.current_language == "tr" else "Warning", warning_msg)
            return
        
        # Sistem kontrolü
        thread_count = self.spin_threads.value()
        
        if PSUTIL_AVAILABLE and thread_count > 16:
            try:
                cpu_count = psutil.cpu_count(logical=True) or 2
                if thread_count > cpu_count * 4:
                    if self.current_language == "tr":
                        message = (f"Seçtiğiniz thread sayısı ({thread_count}) sisteminizin "
                                  f"kapasitesinden ({cpu_count} thread) çok yüksek.\n\n"
                                  f"Bu, performansı DÜŞÜREBİLİR ve sistemi yavaşlatabilir.\n\n"
                                  f"Önerilen: {min(16, cpu_count * 2)} thread\n"
                                  f"Devam etmek istiyor musunuz?")
                        title = "Yüksek Thread Uyarısı"
                        yes_text = "Evet, Devam Et"
                        no_text = "Hayır, Ayarları Düzelt"
                    else:
                        message = (f"Selected thread count ({thread_count}) is much higher than "
                                  f"your system capacity ({cpu_count} threads).\n\n"
                                  f"This may DECREASE performance and slow down the system.\n\n"
                                  f"Recommended: {min(16, cpu_count * 2)} threads\n"
                                  f"Do you want to continue?")
                        title = "High Thread Warning"
                        yes_text = "Yes, Continue"
                        no_text = "No, Fix Settings"
                    
                    reply = QMessageBox.question(
                        self, 
                        title,
                        message,
                        QMessageBox.Yes | QMessageBox.No,
                        QMessageBox.No
                    )
                    
                    if reply == QMessageBox.No:
                        self.spin_threads.setValue(min(16, cpu_count * 2))
                        self.tab_widget.setCurrentIndex(3)  # Ayarlar sekmesine git
                        return
            except Exception as e:
                logging.error(f"Sistem kontrol hatası: {e}")
        
        if self.cb_backup.isChecked():
            self.create_backup()
        
        self.worker_thread = OrganizerWorker()
        self.worker_thread.setup(
            config=self.config,
            ai_modu=ai_mode,
            document_ai_modu=document_ai,
            include_subfolders=self.cb_subfolders.isChecked(),
            delete_empty=self.cb_empty_folders.isChecked(),
            source_folder=self.selected_folder,
            thread_count=thread_count
        )
        
        self.worker_thread.signals.progress.connect(self.progress_bar.setValue)
        self.worker_thread.signals.status.connect(self.status_label.setText)
        self.worker_thread.signals.finished.connect(self.operation_finished)
        self.worker_thread.signals.error.connect(lambda msg: QMessageBox.critical(self, "Hata" if self.current_language == "tr" else "Error", msg))
        self.worker_thread.signals.stats_update.connect(self.update_statistics)
        self.worker_thread.signals.file_processed.connect(self.update_file_list)
        self.worker_thread.signals.ai_analysis_result.connect(self.update_ai_results)
        self.worker_thread.signals.system_info.connect(self.system_info_label.setText)
        self.worker_thread.signals.document_analyzed.connect(self.update_document_analysis)
        
        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        self.btn_start.setEnabled(False)
        self.btn_ai.setEnabled(False)
        self.btn_document_ai.setEnabled(False)
        self.btn_undo.setEnabled(False)
        self.file_list.clear()
        self.ai_results_text.clear()
        self.system_info_label.setText("")
        
        if self.current_language == "tr":
            mode_text = "AI" if ai_mode else ("Doküman AI" if document_ai else "Normal")
            status_text = f"İşlem başlatılıyor ({thread_count} thread, {mode_text})..."
        else:
            mode_text = "AI" if ai_mode else ("Document AI" if document_ai else "Normal")
            status_text = f"Starting process ({thread_count} threads, {mode_text})..."
        
        self.status_label.setText(status_text)
        
        self.worker_thread.start()
    
    def analyze_documents(self):
        """Sadece doküman analizi yap"""
        if not self.selected_folder:
            QMessageBox.warning(self, "Uyarı", "Lütfen önce bir klasör seçin!")
            return
        
        # Doküman dosyalarını topla
        doc_files = []
        doc_extensions = ['.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.txt', '.md']
        
        for root, dirs, files in os.walk(self.selected_folder):
            for filename in files:
                ext = os.path.splitext(filename)[1].lower()
                if ext in doc_extensions:
                    filepath = os.path.join(root, filename)
                    doc_files.append(filepath)
        
        if not doc_files:
            QMessageBox.information(self, "Bilgi", "Analiz edilecek doküman bulunamadı!")
            return
        
        self.document_worker = DocumentAnalysisWorker()
        self.document_worker.setup(doc_files, self.config)
        
        self.document_worker.signals.progress.connect(self.progress_bar.setValue)
        self.document_worker.signals.status.connect(self.status_label.setText)
        self.document_worker.signals.finished.connect(self.document_analysis_finished)
        self.document_worker.signals.error.connect(lambda msg: QMessageBox.critical(self, "Hata", msg))
        self.document_worker.signals.document_analyzed.connect(self.update_document_tree)
        
        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        self.status_label.setText(f"{len(doc_files)} doküman analiz ediliyor...")
        self.document_results = []
        
        self.document_worker.start()
    
    def document_analysis_finished(self, success: bool, result: dict):
        """Doküman analizi tamamlandığında"""
        self.progress_bar.setVisible(False)
        
        if success:
            self.document_results = result.get("results", [])
            total = result.get("total_analyzed", 0)
            
            # İstatistikleri güncelle
            self.update_document_stats()
            
            self.status_label.setText(f"✓ {total} doküman analiz edildi!")
            self.btn_export_analysis.setEnabled(len(self.document_results) > 0)
            
            # Doküman sekmesine geç
            self.tab_widget.setCurrentIndex(1)
        else:
            self.status_label.setText("✗ Doküman analizinde hata oluştu!")
    
    def update_document_tree(self, filename: str, category: str, summary: str, metadata: str):
        """Doküman ağacını güncelle"""
        try:
            # Metadata'dan boyutu al
            meta_dict = json.loads(metadata)
            size = meta_dict.get("dosya_boyutu", 0)
            
            # Boyutu formatla
            for unit in ['B', 'KB', 'MB', 'GB']:
                if size < 1024.0:
                    size_str = f"{size:.1f} {unit}"
                    break
                size /= 1024.0
            else:
                size_str = f"{size:.1f} GB"
            
            item = QTreeWidgetItem([filename, category, summary[:100] + "..." if len(summary) > 100 else summary, size_str])
            self.doc_tree.addTopLevelItem(item)
            
        except Exception as e:
            logging.error(f"Doküman ağacı güncelleme hatası: {e}")
    
    def update_document_stats(self):
        """Doküman istatistiklerini güncelle"""
        if not self.document_results:
            self.doc_stats_text.setText("Henüz analiz yapılmadı.")
            return
        
        # Kategori dağılımı
        categories = Counter([r["category"] for r in self.document_results])
        
        stats_text = f"Toplam Doküman: {len(self.document_results)}\n\n"
        stats_text += "Kategori Dağılımı:\n"
        
        for category, count in categories.most_common():
            percentage = (count / len(self.document_results)) * 100
            stats_text += f"  {category}: {count} (%{percentage:.1f})\n"
        
        self.doc_stats_text.setText(stats_text)
    
    def filter_documents(self):
        """Dokümanları filtrele"""
        filter_type = self.doc_filter_combo.currentText()
        
        for i in range(self.doc_tree.topLevelItemCount()):
            item = self.doc_tree.topLevelItem(i)
            filename = item.text(0)
            ext = os.path.splitext(filename)[1].lower()
            
            show = False
            if filter_type == "Tümü":
                show = True
            elif filter_type == "PDF" and ext == '.pdf':
                show = True
            elif filter_type == "DOCX" and ext == '.docx':
                show = True
            elif filter_type == "PPTX" and ext == '.pptx':
                show = True
            elif filter_type == "Excel" and ext in ['.xlsx', '.xls']:
                show = True
            elif filter_type == "TXT" and ext == '.txt':
                show = True
            
            item.setHidden(not show)
    
    def view_document_details(self, item, column):
        """Doküman detaylarını göster"""
        filename = item.text(0)
        
        # İlgili dokümanı bul
        doc_info = None
        for result in self.document_results:
            if result["filename"] == filename:
                doc_info = result
                break
        
        if doc_info:
            self.doc_summary_text.setText(doc_info["summary"])
            self.doc_metadata_text.setText(json.dumps(doc_info["metadata"], indent=2, ensure_ascii=False))
            self.doc_preview_text.setText(doc_info["text_preview"])
            self.btn_view_doc.setEnabled(True)
    
    def open_document(self):
        """Dokümanı aç"""
        current_item = self.doc_tree.currentItem()
        if not current_item:
            return
        
        filename = current_item.text(0)
        
        # Dosya yolunu bul
        for result in self.document_results:
            if result["filename"] == filename:
                filepath = result["path"]
                QDesktopServices.openUrl(QUrl.fromLocalFile(filepath))
                break
    
    def export_analysis(self):
        """Analiz sonuçlarını dışa aktar"""
        if not self.document_results:
            return
        
        file_path, _ = QFileDialog.getSaveFileName(
            self,
            "Analizi Kaydet",
            f"dokuman_analizi_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
            "JSON Files (*.json);;CSV Files (*.csv)"
        )
        
        if file_path:
            try:
                if file_path.endswith('.json'):
                    with open(file_path, 'w', encoding='utf-8') as f:
                        json.dump(self.document_results, f, indent=4, ensure_ascii=False)
                elif file_path.endswith('.csv'):
                    # CSV için düzleştirilmiş veri
                    import csv
                    with open(file_path, 'w', encoding='utf-8', newline='') as f:
                        writer = csv.writer(f)
                        writer.writerow(['Dosya', 'Kategori', 'Özet', 'Boyut', 'Anahtar Kelimeler'])
                        
                        for result in self.document_results:
                            keywords = ', '.join(result['metadata'].get('anahtar_kelimeler', []))
                            writer.writerow([
                                result['filename'],
                                result['category'],
                                result['summary'],
                                result['metadata'].get('dosya_boyutu', 0),
                                keywords
                            ])
                
                QMessageBox.information(self, "Başarılı", f"Analiz sonuçları kaydedildi: {file_path}")
                
            except Exception as e:
                QMessageBox.critical(self, "Hata", f"Dosya kaydetme hatası: {str(e)}")
    
    def update_document_analysis(self, filename: str, category: str, summary: str, metadata: str):
        """Doküman analizi güncellemesi"""
        # Ana sekmedeki AI sonuçlarına ekle
        current_text = self.ai_results_text.toPlainText()
        new_entry = f"📄 {filename}: {category} - {summary[:50]}...\n"
        
        lines = current_text.split('\n')
        if len(lines) > 10:
            lines = lines[-10:]
            current_text = '\n'.join(lines)
        
        self.ai_results_text.setText(current_text + new_entry)
    
    def operation_finished(self, success: bool, stats: dict):
        """İşlem tamamlandığında"""
        self.progress_bar.setVisible(False)
        self.btn_start.setEnabled(True)
        self.btn_ai.setEnabled(True)
        self.btn_document_ai.setEnabled(True)
        self.check_log_file()
        
        if success:
            requested_threads = stats.get("requested_threads", stats.get("thread_count_used", 1))
            actual_threads = stats.get("thread_count_used", 1)
            
            if requested_threads != actual_threads and self.cb_system_optimization.isChecked():
                if self.current_language == "tr":
                    thread_info = f" (İstenen: {requested_threads}, Kullanılan: {actual_threads} thread)"
                else:
                    thread_info = f" (Requested: {requested_threads}, Used: {actual_threads} threads)"
            else:
                if self.current_language == "tr":
                    thread_info = f" ({actual_threads} thread)"
                else:
                    thread_info = f" ({actual_threads} threads)"
            
            if self.current_language == "tr":
                status_text = f"✓ İşlem başarıyla tamamlandı!{thread_info}"
            else:
                status_text = f"✓ Operation completed successfully!{thread_info}"
            
            self.status_label.setText(status_text)
            self.update_statistics_display(stats)
            self.logger.info(f"İşlem tamamlandı: {stats}")
            self.tab_widget.setCurrentIndex(2)  # İstatistikler sekmesine git
        else:
            if self.current_language == "tr":
                self.status_label.setText("✗ İşlem sırasında hata oluştu!")
            else:
                self.status_label.setText("✗ An error occurred during operation!")
    
    def check_log_file(self):
        """Log dosyasını kontrol et"""
        if not self.selected_folder:
            return
        
        log_path = os.path.join(self.selected_folder, self.log_file)
        if os.path.exists(log_path):
            try:
                with open(log_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    count = len(data) if isinstance(data, list) else 0
                    undo_text = self.translator.get_text('undo', self.current_language)
                    self.btn_undo.setText(f"{undo_text} ({count})")
                    self.btn_undo.setEnabled(count > 0)
            except:
                self.btn_undo.setEnabled(False)
        else:
            self.btn_undo.setText(self.translator.get_text("undo", self.current_language))
            self.btn_undo.setEnabled(False)
    
    def undo_action(self):
        """Son işlemi geri al"""
        log_path = os.path.join(self.selected_folder, self.log_file)
        if not os.path.exists(log_path):
            return
        
        try:
            with open(log_path, 'r', encoding='utf-8') as f:
                history = json.load(f)
            
            if not history:
                return
            
            last_operation = history.pop()
            movements = last_operation.get("movements", [])
            success_count = 0
            
            for move in movements:
                try:
                    if os.path.exists(move["new_path"]):
                        old_dir = os.path.dirname(move["old_path"])
                        os.makedirs(old_dir, exist_ok=True)
                        shutil.move(move["new_path"], move["old_path"])
                        success_count += 1
                except Exception as e:
                    logging.error(f"Geri alma hatası {move['new_path']}: {e}")
            
            if history:
                with open(log_path, 'w', encoding='utf-8') as f:
                    json.dump(history, f, indent=4)
            else:
                os.remove(log_path)
            
            if self.current_language == "tr":
                success_msg = f"✓ {success_count} dosya geri alındı!"
            else:
                success_msg = f"✓ {success_count} files restored!"
            self.status_label.setText(success_msg)
            self.check_log_file()
            self.file_list.clear()
            self.ai_results_text.clear()
            
        except Exception as e:
            if self.current_language == "tr":
                error_msg = f"Geri alma hatası: {str(e)}"
            else:
                error_msg = f"Restore error: {str(e)}"
            QMessageBox.critical(self, "Hata" if self.current_language == "tr" else "Error", error_msg)
            self.logger.error(f"Geri alma hatası: {e}")
    
    def create_backup(self):
        """Yedek oluştur"""
        try:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            backup_dir = os.path.join(self.selected_folder, f"backup_{timestamp}")
            
            if not os.path.exists(backup_dir):
                os.makedirs(backup_dir)
                
                for item in os.listdir(self.selected_folder):
                    item_path = os.path.join(self.selected_folder, item)
                    if os.path.isfile(item_path) and not item.startswith(".") and not item.startswith("backup_"):
                        try:
                            shutil.copy2(item_path, os.path.join(backup_dir, item))
                        except Exception as e:
                            logging.warning(f"Yedekleme sırasında {item} kopyalanamadı: {e}")
                
                self.logger.info(f"Yedek oluşturuldu: {backup_dir}")
                
        except Exception as e:
            self.logger.error(f"Yedek oluşturma hatası: {e}")
    
    def update_statistics(self, stats: dict):
        """İstatistikleri güncelle"""
        pass
    
    def update_statistics_display(self, stats: dict):
        """İstatistikleri görüntüle"""
        self.stats_table.setRowCount(8)
        
        if self.current_language == "tr":
            rows = [
                ["Toplam Dosya", str(stats.get("total_files", 0))],
                ["İşlenen Dosya", str(stats.get("processed_files", 0))],
                ["Atlanan Dosya", str(stats.get("skipped_files", 0))],
                ["Oluşturulan Klasör", str(stats.get("created_folders", 0))],
                ["İstenen Thread Sayısı", str(stats.get("requested_threads", stats.get("thread_count_used", 1)))],
                ["Kullanılan Thread Sayısı", str(stats.get("thread_count_used", 1))],
                ["Toplam Boyut", stats.get("formatted_size", "0 B")],
                ["İşlem Süresi", f"{stats.get('processing_time', 0):.2f} saniye"]
            ]
        else:
            rows = [
                ["Total Files", str(stats.get("total_files", 0))],
                ["Processed Files", str(stats.get("processed_files", 0))],
                ["Skipped Files", str(stats.get("skipped_files", 0))],
                ["Created Folders", str(stats.get("created_folders", 0))],
                ["Requested Threads", str(stats.get("requested_threads", stats.get("thread_count_used", 1)))],
                ["Threads Used", str(stats.get("thread_count_used", 1))],
                ["Total Size", stats.get("formatted_size", "0 B")],
                ["Processing Time", f"{stats.get('processing_time', 0):.2f} seconds"]
            ]
        
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                item = QTableWidgetItem(cell)
                item.setTextAlignment(Qt.AlignCenter)
                self.stats_table.setItem(i, j, item)
        
        # Kategori dağılımı
        categories = stats.get("category_distribution", {})
        self.category_table.setRowCount(len(categories))
        
        if self.current_language == "tr":
            self.category_table.setHorizontalHeaderLabels(["Kategori", "Dosya Sayısı"])
        else:
            self.category_table.setHorizontalHeaderLabels(["Category", "File Count"])
        
        for i, (category, count) in enumerate(categories.items()):
            category_item = QTableWidgetItem(category)
            category_item.setTextAlignment(Qt.AlignCenter)
            self.category_table.setItem(i, 0, category_item)
            
            count_item = QTableWidgetItem(str(count))
            count_item.setTextAlignment(Qt.AlignCenter)
            self.category_table.setItem(i, 1, count_item)
        
        # Doküman kategorileri
        doc_categories = stats.get("document_categories", {})
        self.document_category_table.setRowCount(len(doc_categories))
        
        if self.current_language == "tr":
            self.document_category_table.setHorizontalHeaderLabels(["Kategori", "Doküman Sayısı"])
        else:
            self.document_category_table.setHorizontalHeaderLabels(["Category", "Document Count"])
        
        for i, (category, count) in enumerate(doc_categories.items()):
            category_item = QTableWidgetItem(category)
            category_item.setTextAlignment(Qt.AlignCenter)
            self.document_category_table.setItem(i, 0, category_item)
            
            count_item = QTableWidgetItem(str(count))
            count_item.setTextAlignment(Qt.AlignCenter)
            self.document_category_table.setItem(i, 1, count_item)
        
        # Nesneler
        objects = stats.get("detected_objects", {})
        self.objects_table.setRowCount(len(objects))
        
        if self.current_language == "tr":
            self.objects_table.setHorizontalHeaderLabels(["Nesne", "Sayı"])
        else:
            self.objects_table.setHorizontalHeaderLabels(["Object", "Count"])
        
        for i, (obj, count) in enumerate(objects.items()):
            obj_item = QTableWidgetItem(obj)
            obj_item.setTextAlignment(Qt.AlignCenter)
            self.objects_table.setItem(i, 0, obj_item)
            
            count_item = QTableWidgetItem(str(count))
            count_item.setTextAlignment(Qt.AlignCenter)
            self.objects_table.setItem(i, 1, count_item)
    
    def update_file_list(self, filename: str, category: str):
        """Dosya listesini güncelle"""
        current_text = self.file_list.toPlainText()
        new_entry = f"✓ {filename} → {category}\n"
        self.file_list.setText(current_text + new_entry)
        
        # DÜZELTME: cursor.movePosition'u düzelt
        cursor = self.file_list.textCursor()
        cursor.movePosition(QTextCursor.MoveOperation.End)
        self.file_list.setTextCursor(cursor)
    
    def update_ai_results(self, filename: str, category: str, objects: List[str]):
        """AI analiz sonuçlarını güncelle"""
        if objects:
            objects_str = ", ".join(objects[:3])
            current_text = self.ai_results_text.toPlainText()
            
            if self.current_language == "tr":
                new_entry = f"📊 {filename}: {category} [Nesneler: {objects_str}]\n"
            else:
                new_entry = f"📊 {filename}: {category} [Objects: {objects_str}]\n"
            
            lines = current_text.split('\n')
            if len(lines) > 10:
                lines = lines[-10:]
                current_text = '\n'.join(lines)
            
            self.ai_results_text.setText(current_text + new_entry)
            
            # DÜZELTME: cursor.movePosition'u düzelt
            cursor = self.ai_results_text.textCursor()
            cursor.movePosition(QTextCursor.MoveOperation.End)
            self.ai_results_text.setTextCursor(cursor)
    
    def change_language(self, language: str):
        """Dili değiştir"""
        if language == "Türkçe":
            self.current_language = "tr"
        else:
            self.current_language = "en"
        
        self.config["language"] = self.current_language
        self.update_ui_texts()
    
    def change_theme(self, theme: str):
        """Temayı değiştir"""
        if theme == "Açık":
            self.current_theme = "light"
        else:
            self.current_theme = "dark"
        
        self.config["theme"] = self.current_theme
        self.apply_theme()
        self.config_manager.save_config(self.config)
    
    def save_settings(self):
        """Ayarları kaydet"""
        self.config["ai_enabled"] = self.cb_ai_enabled.isChecked()
        self.config["document_ai_enabled"] = self.cb_document_ai_enabled.isChecked()
        self.config["thread_count"] = self.spin_threads.value()
        self.config["backup_enabled"] = self.cb_backup.isChecked()
        self.config["multi_object_detection"] = self.cb_multi_object.isChecked()
        self.config["object_threshold"] = self.threshold_slider.value()
        self.config["max_objects"] = self.spin_max_objects.value()
        self.config["enable_system_optimization"] = self.cb_system_optimization.isChecked()
        
        # Analiz derinliği
        depth_index = self.analysis_depth_combo.currentIndex()
        depth_map = {0: "low", 1: "medium", 2: "high"}
        self.config["document_analysis_depth"] = depth_map.get(depth_index, "medium")
        
        self.config_manager.save_config(self.config)
        success_msg = "Ayarlar kaydedildi!" if self.current_language == "tr" else "Settings saved!"
        QMessageBox.information(self, "Başarılı" if self.current_language == "tr" else "Success", success_msg)
    
    def update_ui_texts(self):
        """UI metinlerini güncelle"""
        self.setWindowTitle(self.translator.get_text("app_title", self.current_language))
        
        # Başlık
        title_text = "📁 AI Doküman Düzenleyici" if self.current_language == "tr" else "📁 AI Document Organizer"
        self.title_label.setText(title_text)
        
        # Buton metinleri
        self.btn_select.setText(self.translator.get_text("select_folder", self.current_language))
        self.btn_view.setText(self.translator.get_text("view_folder", self.current_language))
        self.btn_start.setText(self.translator.get_text("start", self.current_language))
        self.btn_ai.setText(self.translator.get_text("ai_mode", self.current_language))
        self.btn_document_ai.setText(self.translator.get_text("document_ai_mode", self.current_language))
        self.btn_undo.setText(self.translator.get_text("undo", self.current_language))
        self.btn_save_settings.setText(self.translator.get_text("save_settings", self.current_language))
        self.btn_analyze_documents.setText(self.translator.get_text("analyze_documents", self.current_language))
        self.btn_view_doc.setText(self.translator.get_text("view_document", self.current_language))
        self.btn_export_analysis.setText(self.translator.get_text("export_analysis", self.current_language))
        
        # Sekme metinleri
        self.tab_widget.setTabText(0, self.translator.get_text("main", self.current_language))
        self.tab_widget.setTabText(1, self.translator.get_text("documents", self.current_language))
        self.tab_widget.setTabText(2, self.translator.get_text("stats", self.current_language))
        self.tab_widget.setTabText(3, self.translator.get_text("settings", self.current_language))
        
        # CheckBox metinleri
        self.cb_subfolders.setText(self.translator.get_text("scan_subfolders", self.current_language))
        self.cb_empty_folders.setText(self.translator.get_text("delete_empty", self.current_language))
        self.cb_backup.setText(self.translator.get_text("create_backup", self.current_language))
        self.cb_ai_enabled.setText(self.translator.get_text("ai_enabled", self.current_language))
        self.cb_document_ai_enabled.setText(self.translator.get_text("document_ai_enabled", self.current_language))
        self.cb_multi_object.setText(self.translator.get_text("multi_object", self.current_language))
        self.cb_system_optimization.setText(self.translator.get_text("system_optimization", self.current_language))
        self.cb_system_optimization.setToolTip(self.translator.get_text("system_optimization_tooltip", self.current_language))
        
        # GroupBox başlıkları
        self.folder_group.setTitle(self.translator.get_text("folder_selection", self.current_language))
        self.options_group.setTitle(self.translator.get_text("options", self.current_language))
        self.ai_group.setTitle(self.translator.get_text("ai_results", self.current_language))
        self.file_group.setTitle(self.translator.get_text("file_history", self.current_language))
        self.general_group.setTitle(self.translator.get_text("general_settings", self.current_language))
        self.ai_settings_group.setTitle(self.translator.get_text("ai_settings", self.current_language))
        self.document_settings_group.setTitle(self.translator.get_text("document_settings", self.current_language))
        
        # Statik GroupBox başlıkları
        if self.current_language == "tr":
            self.stats_group.setTitle("İstatistikler")
            self.category_group.setTitle("Kategori Dağılımı")
            self.document_category_group.setTitle("Doküman Kategorileri")
            self.objects_group.setTitle("Tespit Edilen Nesneler")
        else:
            self.stats_group.setTitle("Statistics")
            self.category_group.setTitle("Category Distribution")
            self.document_category_group.setTitle("Document Categories")
            self.objects_group.setTitle("Detected Objects")
        
        # Etiket metinleri
        self.lang_label.setText(f"{self.translator.get_text('language', self.current_language)}:")
        self.theme_label.setText(f"{self.translator.get_text('theme', self.current_language)}:")
        self.thread_label.setText(f"{self.translator.get_text('thread_count', self.current_language)}:")
        self.threshold_label.setText(f"{self.translator.get_text('object_threshold', self.current_language)}:")
        self.max_objects_label.setText(f"{self.translator.get_text('max_objects', self.current_language)}:")
        
        # Thread bilgi metni
        self.thread_info_label.setText(self.translator.get_text("thread_info", self.current_language))
        
        # Placeholder metinleri
        if self.current_language == "tr":
            placeholder_text = "AI analiz sonuçları burada gösterilecek..."
            folder_text = "Henüz klasör seçilmedi"
            doc_summary_placeholder = "Doküman özeti burada gösterilecek..."
            doc_metadata_placeholder = "Metadata bilgileri burada gösterilecek..."
            doc_preview_placeholder = "Metin önizlemesi burada gösterilecek..."
            doc_stats_placeholder = "Analiz istatistikleri burada gösterilecek..."
        else:
            placeholder_text = "AI analysis results will be shown here..."
            folder_text = "No folder selected yet"
            doc_summary_placeholder = "Document summary will be shown here..."
            doc_metadata_placeholder = "Metadata information will be shown here..."
            doc_preview_placeholder = "Text preview will be shown here..."
            doc_stats_placeholder = "Analysis statistics will be shown here..."
        
        self.ai_results_text.setPlaceholderText(placeholder_text)
        self.folder_label.setText(folder_text)
        self.doc_summary_text.setPlaceholderText(doc_summary_placeholder)
        self.doc_metadata_text.setPlaceholderText(doc_metadata_placeholder)
        self.doc_preview_text.setPlaceholderText(doc_preview_placeholder)
        self.doc_stats_text.setPlaceholderText(doc_stats_placeholder)
        
        # Statü çubuğu
        self.statusBar().showMessage("Hazır" if self.current_language == "tr" else "Ready")

# ==================== APPLICATION ENTRY ====================

def main():
    """Uygulama giriş noktası"""
    multiprocessing.freeze_support()
    
    if hasattr(Qt, 'AA_EnableHighDpiScaling'):
        QApplication.setAttribute(Qt.AA_EnableHighDpiScaling, True)
    if hasattr(Qt, 'AA_UseHighDpiPixmaps'):
        QApplication.setAttribute(Qt.AA_UseHighDpiPixmaps, True)
    
    app = QApplication(sys.argv)
    app.setApplicationName("AI Doküman Düzenleyici")
    app.setOrganizationName("FileOrganizer")
    
    window = DosyaDuzenleyici()
    window.show()
    
    sys.exit(app.exec())

if __name__ == "__main__":
    main()